#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
禾灵儿文化德育数字画像系统 —— 核心引擎与 Web 应用（单文件自包含）
===============================================================================

用途
    1) 直接作为"禾灵儿"文化智能体 2.0 网站的一个栏目部署：
           uvicorn hlr_portrait:app --host 0.0.0.0 --port 8080
       或在既有 FastAPI 工程中挂载：
           from hlr_portrait import app as portrait_app
           main_app.mount("/portrait", portrait_app)
    2) 被 desktop_app.py 引用，打包成 Windows 单文件 exe 供教师本地使用。

评价框架
    五个维度：文化体验 / 文化认知 / 文化志趣 / 文化认同 / 文化研创
    取值区间 [0, 1]，满分 1.0。

    初始得分来自开课第一节课的五维问卷；此后学生在平台上的每一次学习、作业、
    创作、智能体开发都会转化为"增量证据"，令对应维度得分单调上升并趋近 1.0，
    但永远不会超过 1.0：

        S_d = B_d + (1 - B_d) * (1 - exp(-A_d / kappa_d))

    其中 B_d 为初始得分，A_d 为该维度累计的有效证据量。该式保证：
        · A_d = 0        -> S_d = B_d          （只做问卷时就是初始得分）
        · A_d 增大       -> S_d 严格单调上升    （材料越多分越高）
        · A_d -> 无穷    -> S_d -> 1.0          （封顶 1.0，永不越界）

确定性与稳定性
    全流程不调用任何大模型、不依赖网络，打分只由"词典 + 规则 + 计量指标"决定。
    每条证据都带内容指纹（含文件 SHA-256），重复上传同一材料只计一次。
    因此"同样的内容反复上传，得分完全稳定"是由存储层强制保证的，而非近似。

作者信息与许可见 README.md。
"""
from __future__ import annotations

import hashlib
import io
import json
import math
import os
import re
import shutil
import sqlite3
import struct
import sys
import tempfile
import threading
import unicodedata
from dataclasses import dataclass, field
from datetime import datetime, timezone
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

# --------------------------------------------------------------------------
# 可选依赖：缺失时对应格式自动降级，不影响其余功能
# --------------------------------------------------------------------------
# 说明：这里刻意捕获 BaseException 而非 Exception。个别环境中损坏的 cryptography
# 会让 pypdf 的导入抛出 pyo3 的 PanicException（不继承自 Exception），
# 若只捕获 Exception 会导致整个程序无法启动。
try:
    import openpyxl
except BaseException:                                # pragma: no cover
    openpyxl = None

try:
    import docx as _docx                             # python-docx
except BaseException:                                # pragma: no cover
    _docx = None

try:
    from pptx import Presentation as _Presentation   # python-pptx
except BaseException:                                # pragma: no cover
    _Presentation = None

# pypdf 延迟到真正遇到 PDF 时再导入：部分环境里它的加密后端会在导入期
# 打印 Rust panic 噪声，延迟导入可让程序启动保持干净。
_PdfReader = None
_PDF_TRIED = False


def _pdf_reader():                                   # pragma: no cover
    global _PdfReader, _PDF_TRIED
    if not _PDF_TRIED:
        _PDF_TRIED = True
        try:
            from pypdf import PdfReader
            _PdfReader = PdfReader
        except BaseException:
            _PdfReader = None
    return _PdfReader


# ==========================================================================
# 一、常量与可配置项
# ==========================================================================
APP_NAME = "禾灵儿文化德育数字画像"
APP_VERSION = "1.0.0"

DIMS: Tuple[str, ...] = ("文化体验", "文化认知", "文化志趣", "文化认同", "文化研创")
DIM_CODE: Dict[str, str] = {
    "文化体验": "E", "文化认知": "C", "文化志趣": "I",
    "文化认同": "D", "文化研创": "R",
}
CODE_DIM: Dict[str, str] = {v: k for k, v in DIM_CODE.items()}

# 界面上的年级列（对应截图中的"任务1~任务4"位置）
GRADES: Tuple[str, ...] = ("大一", "大二", "大三", "大四")
GRADE_INDEX: Dict[str, int] = {g: i for i, g in enumerate(GRADES)}

BASELINE_DEFAULT = 0.50          # 未做问卷时的中性初始得分
EPS = 1e-9

# 各维度的"证据饱和常数"：A_d 达到 kappa 时，与满分的差距缩小约 63%。
# 数值越小，同样的材料带来的提分越快。研创类材料稀缺，故取值最小。
KAPPA: Dict[str, float] = {
    "文化体验": 4.5,
    "文化认知": 5.0,
    "文化志趣": 4.5,
    "文化认同": 4.0,
    "文化研创": 3.5,
}

# 单一渠道对某维度的证据量上限。在汇总时截断（而非按上传批次截断），
# 因此与上传顺序无关，重算结果始终一致，同时抑制"单一渠道刷量"。
CHANNEL_A_CAP: Dict[str, float] = {
    "dialogue": 6.0, "enrollment": 6.0, "homework": 8.0,
    "video": 6.0, "artifact": 12.0, "agent": 12.0,
}

# 渠道 -> 维度 贡献权重（行内不必归一，直接作为乘子）
CHANNEL_WEIGHTS: Dict[str, Dict[str, float]] = {
    # 与禾灵儿智能体的文化对话
    "dialogue":   {"文化体验": 0.20, "文化认知": 0.55, "文化志趣": 0.50,
                   "文化认同": 0.40, "文化研创": 0.25},
    # 文化课程报名（广度与主动选学）
    "enrollment": {"文化体验": 0.40, "文化认知": 0.25, "文化志趣": 0.60,
                   "文化认同": 0.25, "文化研创": 0.10},
    # 课程作业完成与得分
    "homework":   {"文化体验": 0.25, "文化认知": 0.65, "文化志趣": 0.40,
                   "文化认同": 0.20, "文化研创": 0.35},
    # 课程视频学习完成率
    "video":      {"文化体验": 0.60, "文化认知": 0.45, "文化志趣": 0.35,
                   "文化认同": 0.20, "文化研创": 0.05},
    # 学生提交的作业成果（文档 / 演示 / 表格 / 视频）
    "artifact":   {"文化体验": 0.35, "文化认知": 0.45, "文化志趣": 0.30,
                   "文化认同": 0.40, "文化研创": 0.85},
    # 学生自建的文化智能体
    "agent":      {"文化体验": 0.20, "文化认知": 0.50, "文化志趣": 0.40,
                   "文化认同": 0.35, "文化研创": 1.00},
}

CHANNEL_LABEL: Dict[str, str] = {
    "questionnaire": "五维问卷（初始得分）",
    "dialogue": "文化对话",
    "enrollment": "文化课程报名",
    "homework": "课程作业",
    "video": "课程视频",
    "artifact": "学生作业成果",
    "agent": "学生自建智能体",
}

# --------------------------------------------------------------------------
# 文化主题词典
#   GENERIC   通用文化词      REGION 嘉兴/江南地域文化词      REVOLUTION 革命文化词
#   DIM_CUES  五个维度的语义线索词
# 词典可由同目录 lexicon.json 覆盖，便于教师后期扩充而无需改代码。
# --------------------------------------------------------------------------
LEX_GENERIC = [
    "文化", "传统", "非遗", "非物质文化遗产", "民俗", "历史", "文脉", "艺术",
    "美学", "传承", "遗产", "古镇", "民族", "中华", "国学", "经典", "诗词",
    "书法", "绘画", "戏曲", "民歌", "手工艺", "技艺", "习俗", "节气", "节日",
    "古建筑", "文物", "考古", "典籍", "礼仪", "家风", "乡土", "地方志",
    "博物馆", "纪念馆", "美术馆", "文创", "国潮", "汉服", "茶文化", "陶瓷",
    "剪纸", "年画", "皮影", "刺绣", "织造", "木雕", "石刻", "碑刻", "楹联",
]
LEX_REGION = [
    "嘉兴", "禾城", "南湖", "乌镇", "西塘", "南浔", "月河", "运河", "京杭大运河",
    "江南", "水乡", "浙北", "桐乡", "平湖", "海盐", "海宁", "秀洲", "嘉善",
    "粽子", "五芳斋", "蓝印花布", "灶画", "掼牛", "皮影", "硖石灯彩", "田歌",
    "船拳", "网船会", "三跳", "紫砂", "浔溪", "子城", "文生修船", "槜李",
    "茅盾", "丰子恺", "金庸", "王国维", "徐志摩", "李叔同", "沈钧儒", "张乐平",
    "马家浜", "良渚", "崧泽", "运河文化带", "江南文化", "水乡文脉",
]
LEX_REVOLUTION = [
    "红船", "红船精神", "建党", "一大", "革命", "烈士", "长征", "抗战", "初心",
    "使命", "党史", "南湖革命纪念馆", "红色", "先烈", "延安", "井冈山", "五四",
    "爱国", "统一战线", "社会主义先进文化", "中国共产党",
]
LEX_DIM_CUES: Dict[str, List[str]] = {
    "文化体验": [
        "参观", "走访", "研学", "实地", "打卡", "漫游", "体验", "现场", "亲历",
        "沉浸", "展馆", "场馆", "展品", "街区", "景区", "遗址", "游览", "采风",
        "考察", "调研", "手作", "亲手", "参与", "工作坊", "实践活动", "志愿",
        "夏令营", "游学", "踏访", "寻访",
    ],
    "文化认知": [
        "是什么", "为什么", "起源", "由来", "发展", "含义", "内涵", "区别",
        "特点", "特征", "原因", "解释", "知识", "概念", "演变", "历史沿革",
        "分类", "结构", "原理", "背景", "考证", "辨析", "综述", "梳理",
        "怎么形成", "如何理解", "有哪些",
    ],
    # 注意：“推荐/介绍”表达的是传播意愿，归入文化认同而非文化志趣
    "文化志趣": [
        "想了解", "感兴趣", "喜欢", "还有哪些", "继续", "更多", "深入",
        "拓展", "追问", "再讲讲", "详细说说", "接着", "另外", "此外", "延伸",
        "自学", "主动", "坚持", "每天", "持续", "课外", "业余",
    ],
    "文化认同": [
        "我们", "自豪", "骄傲", "家乡", "认同", "归属", "弘扬", "传播",
        "推荐", "介绍", "分享给", "传承",
        "文化自信", "讲好中国故事", "向世界", "对外", "国际", "翻译", "外译",
        "责任", "担当", "使命感", "情怀", "热爱", "根脉", "乡愁", "民族精神",
    ],
    "文化研创": [
        "设计", "创作", "生成", "制作", "方案", "脚本", "剧本", "策划", "文创",
        "作品", "创意", "海报", "视频", "短片", "微视频", "logo", "标识",
        "开发", "智能体", "提示词", "研究", "论文", "调查报告", "课题",
        "原创", "改编", "再创作", "产品", "衍生品", "插画", "动画",
    ],
}

# 低质量对话的常见口水句，直接降权
LEX_FILLER = {
    "你好", "您好", "hi", "hello", "在吗", "继续", "好的", "谢谢", "嗯", "哦",
    "?", "？", "。", "1", "test", "测试", "再来一个", "还有吗", "ok",
}

VIDEO_EXT = {".mp4", ".mov", ".avi", ".mkv", ".wmv", ".flv", ".m4v", ".webm", ".rmvb"}
DOC_EXT = {".docx", ".doc"}
PPT_EXT = {".pptx", ".ppt"}
XLS_EXT = {".xlsx", ".xlsm", ".xls", ".csv"}
PDF_EXT = {".pdf"}
TEXT_EXT = {".txt", ".md", ".json", ".yaml", ".yml", ".log"}
IMAGE_EXT = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"}


def config_dir() -> str:
    """配置文件目录。

    打包成 exe 后 __file__ 指向 PyInstaller 的临时解包目录，那里放的是只读副本；
    教师要改的 lexicon.json / weights.json / questionnaire_map.json 应当放在
    exe 旁边，因此冻结状态下必须返回 exe 所在目录。
    """
    if getattr(sys, "frozen", False):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


def _load_external_config() -> None:
    """允许用同目录下的 lexicon.json / weights.json 覆盖内置默认值。"""
    base = config_dir()
    lex_path = os.path.join(base, "lexicon.json")
    if os.path.isfile(lex_path):
        try:
            with open(lex_path, "r", encoding="utf-8") as fh:
                data = json.load(fh)
            for key, target in (("generic", LEX_GENERIC), ("region", LEX_REGION),
                                ("revolution", LEX_REVOLUTION)):
                if isinstance(data.get(key), list):
                    target[:] = [str(x) for x in data[key]]
            if isinstance(data.get("dim_cues"), dict):
                for d, words in data["dim_cues"].items():
                    if d in LEX_DIM_CUES and isinstance(words, list):
                        LEX_DIM_CUES[d] = [str(x) for x in words]
        except Exception:
            pass
    w_path = os.path.join(base, "weights.json")
    if os.path.isfile(w_path):
        try:
            with open(w_path, "r", encoding="utf-8") as fh:
                data = json.load(fh)
            for ch, mapping in (data.get("channel_weights") or {}).items():
                if ch in CHANNEL_WEIGHTS and isinstance(mapping, dict):
                    CHANNEL_WEIGHTS[ch].update(
                        {k: float(v) for k, v in mapping.items() if k in DIMS})
            for d, v in (data.get("kappa") or {}).items():
                if d in KAPPA:
                    KAPPA[d] = float(v)
        except Exception:
            pass


_load_external_config()


# ==========================================================================
# 二、通用工具
# ==========================================================================
def norm_text(s: Any) -> str:
    """全角转半角、去空白、统一小写，便于稳定匹配。"""
    if s is None:
        return ""
    s = unicodedata.normalize("NFKC", str(s))
    return re.sub(r"\s+", "", s).strip().lower()


def norm_name(s: Any) -> str:
    """人名/昵称规范化：去掉常见的班级前后缀装饰。"""
    t = norm_text(s)
    t = re.sub(r"[\-—_·．.]+", "", t)
    return t


def norm_sid(s: Any) -> str:
    """学号规范化：仅保留数字与字母。"""
    if s is None:
        return ""
    t = unicodedata.normalize("NFKC", str(s)).strip()
    if t.endswith(".0"):
        t = t[:-2]
    return re.sub(r"[^0-9A-Za-z]", "", t)


def parse_percent(v: Any) -> float:
    """'85%' / 0.85 / 85 -> 0.85，异常返回 0.0。"""
    if v is None:
        return 0.0
    if isinstance(v, (int, float)):
        x = float(v)
        return x / 100.0 if x > 1.0 else max(0.0, x)
    t = norm_text(v).replace("%", "")
    if not t:
        return 0.0
    try:
        x = float(t)
    except ValueError:
        return 0.0
    return x / 100.0 if x > 1.0 else max(0.0, x)


def to_int(v: Any, default: int = 0) -> int:
    try:
        if v is None or (isinstance(v, str) and not v.strip()):
            return default
        return int(float(str(v).strip()))
    except Exception:
        return default


def to_float(v: Any, default: float = 0.0) -> float:
    try:
        if v is None or (isinstance(v, str) and not v.strip()):
            return default
        return float(str(v).strip())
    except Exception:
        return default


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def sha256_text(*parts: Any) -> str:
    h = hashlib.sha256()
    for p in parts:
        h.update(str(p).encode("utf-8", "ignore"))
        h.update(b"\x1f")
    return h.hexdigest()


def clamp(x: float, lo: float = 0.0, hi: float = 1.0) -> float:
    return max(lo, min(hi, x))


def now_iso() -> str:
    return datetime.now(timezone.utc).astimezone().strftime("%Y-%m-%d %H:%M:%S")


def normalize_grade(raw: Any) -> Optional[str]:
    """把脏乱的年级字段归一到 大一~大四；无法判定返回 None。

    实测导出数据里出现过：'大一' '2025级' '2025' '25级' '25' '一年级' '1'
    '2025届' '2025级大一' '一' '1年级' '2025级硕士研究生' 等等。
    """
    t = norm_text(raw)
    if not t:
        return None
    for g in GRADES:
        if g in t:
            return g
    if "研究生" in t or "硕士" in t or "博士" in t:
        return None
    cn = {"一": "大一", "二": "大二", "三": "大三", "四": "大四"}
    m = re.search(r"([一二三四])年级", t)
    if m:
        return cn[m.group(1)]
    m = re.fullmatch(r"([一二三四])", t)
    if m:
        return cn[m.group(1)]
    m = re.fullmatch(r"([1-4])(年级)?", t)
    if m:
        return GRADES[int(m.group(1)) - 1]
    # 入学年份形式：2025 / 25 / 2025级 / 25届
    m = re.search(r"(20)?(\d{2})\s*(级|届)?", t)
    if m:
        yy = int(m.group(2))
        year = 2000 + yy
        this_year = datetime.now().year
        month = datetime.now().month
        # 9 月之后算新学年
        academic = this_year - (0 if month >= 9 else 1)
        idx = academic - year
        if 0 <= idx <= 3:
            return GRADES[idx]
    return None


# ==========================================================================
# 三、文本的文化特征分析（确定性，无模型）
# ==========================================================================
@dataclass
class TextSignal:
    """一段文本的文化特征画像。"""
    length: int = 0
    generic: int = 0            # 通用文化词命中数（去重）
    region: int = 0             # 地域文化词命中数（去重）
    revolution: int = 0         # 革命文化词命中数（去重）
    dim_hits: Dict[str, int] = field(default_factory=lambda: {d: 0 for d in DIMS})
    keywords: List[str] = field(default_factory=list)

    @property
    def culture_hits(self) -> int:
        return self.generic + self.region + self.revolution

    @property
    def relevance(self) -> float:
        """文化相关度 [0.30, 1.00]。

        取 0.40 保底是刻意的“稍宽松”：只要学生在平台上认真做事，
        即便这次内容不直接谈文化，也应当获得基础的过程性认可；
        命中地域文化与革命文化词的材料则显著加权。
        """
        score = 0.40 + 0.08 * self.generic + 0.12 * self.region + 0.12 * self.revolution
        return clamp(score, 0.40, 1.00)


def analyze_text(text: str, max_keywords: int = 12) -> TextSignal:
    """统计文本中的文化词与维度线索词。命中按"去重词种数"计，避免复读刷分。"""
    sig = TextSignal()
    if not text:
        return sig
    low = unicodedata.normalize("NFKC", text).lower()
    sig.length = len(low)

    hit_words: List[Tuple[str, str]] = []
    for word in LEX_GENERIC:
        if word.lower() in low:
            sig.generic += 1
            hit_words.append((word, "通用"))
    for word in LEX_REGION:
        if word.lower() in low:
            sig.region += 1
            hit_words.append((word, "地域"))
    for word in LEX_REVOLUTION:
        if word.lower() in low:
            sig.revolution += 1
            hit_words.append((word, "革命"))
    # 志趣类线索词较为泛化（“主动/更多/持续”等在任何题干里都可能出现），
    # 计数时折减 20%，避免把认同类、体验类题目误归到志趣。
    for dim, words in LEX_DIM_CUES.items():
        cnt = sum(1 for word in words if word.lower() in low)
        sig.dim_hits[dim] = cnt

    # 关键词展示优先级：地域 > 革命 > 通用
    order = {"地域": 0, "革命": 1, "通用": 2}
    hit_words.sort(key=lambda x: order[x[1]])
    sig.keywords = [w for w, _ in hit_words[:max_keywords]]
    return sig


def dim_distribution(sig: TextSignal) -> Dict[str, float]:
    """把一段文本的维度线索转成归一化分配比例。

    无任何线索时退化为均分，保证"只要做事就有普遍性收益"。
    """
    raw = {d: float(sig.dim_hits.get(d, 0)) for d in DIMS}
    # 文化词本身对认知/认同有基础贡献
    raw["文化认知"] += 0.5 * sig.generic
    raw["文化认同"] += 0.6 * sig.region + 0.8 * sig.revolution
    total = sum(raw.values())
    if total <= EPS:
        return {d: 1.0 / len(DIMS) for d in DIMS}
    # 平滑：保留 25% 均分，避免单维度独吞
    smooth = 0.25
    return {d: smooth / len(DIMS) + (1 - smooth) * raw[d] / total for d in DIMS}


# ==========================================================================
# 四、文件文本抽取
# ==========================================================================
def read_mp4_duration(path: str) -> Optional[float]:
    """极简 MP4/MOV 时长解析（读取 mvhd box），失败返回 None，不引入第三方依赖。"""
    try:
        with open(path, "rb") as fh:
            data = fh.read(4 * 1024 * 1024)      # mvhd 通常在文件头部
        idx = data.find(b"mvhd")
        if idx < 0:
            return None
        p = idx + 4
        version = data[p]
        p += 4                                    # version(1) + flags(3)
        if version == 1:
            p += 16                               # creation/modification time (8+8)
            timescale = struct.unpack(">I", data[p:p + 4])[0]
            p += 4
            duration = struct.unpack(">Q", data[p:p + 8])[0]
        else:
            p += 8                                # creation/modification time (4+4)
            timescale = struct.unpack(">I", data[p:p + 4])[0]
            p += 4
            duration = struct.unpack(">I", data[p:p + 4])[0]
        if timescale <= 0:
            return None
        return duration / float(timescale)
    except Exception:
        return None


def extract_text_from_file(path: str, filename: str = "") -> Tuple[str, Dict[str, Any]]:
    """从各类作业文件中抽取纯文本与结构性计量指标。

    返回 (text, meta)。meta 里的计量指标参与"研创"维度打分。
    """
    name = filename or os.path.basename(path)
    ext = os.path.splitext(name)[1].lower()
    meta: Dict[str, Any] = {
        "文件名": name, "类型": ext.lstrip(".") or "未知",
        "大小MB": round(os.path.getsize(path) / 1048576.0, 2),
    }
    text = ""

    try:
        if ext in DOC_EXT and _docx is not None and ext == ".docx":
            doc = _docx.Document(path)
            paras = [p.text for p in doc.paragraphs]
            n_tbl = len(doc.tables)
            for tbl in doc.tables:
                for row in tbl.rows:
                    paras.extend(c.text for c in row.cells)
            n_img = sum(1 for r in doc.part.rels.values()
                        if "image" in str(r.reltype))
            text = "\n".join(paras)
            meta.update({"段落数": len(doc.paragraphs), "表格数": n_tbl, "插图数": n_img})

        elif ext in PPT_EXT and _Presentation is not None and ext == ".pptx":
            prs = _Presentation(path)
            chunks, n_img = [], 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
                    if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                        n_img += 1
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    chunks.append(slide.notes_slide.notes_text_frame.text)
            text = "\n".join(chunks)
            meta.update({"页数": len(prs.slides), "插图数": n_img})

        elif ext in XLS_EXT and ext in (".xlsx", ".xlsm") and openpyxl is not None:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            chunks, cells = [], 0
            for ws in wb.worksheets:
                chunks.append(str(ws.title))
                for row in ws.iter_rows(values_only=True):
                    for c in row:
                        if c is not None:
                            chunks.append(str(c))
                            cells += 1
                    if cells > 20000:
                        break
            wb.close()
            text = "\n".join(chunks)
            meta.update({"有效单元格": cells})

        elif ext == ".csv":
            with open(path, "r", encoding="utf-8", errors="ignore") as fh:
                text = fh.read(2 * 1024 * 1024)
            meta.update({"行数": text.count("\n") + 1})

        elif ext in PDF_EXT and _pdf_reader() is not None:
            reader = _pdf_reader()(path)
            text = "\n".join((pg.extract_text() or "") for pg in reader.pages[:80])
            meta.update({"页数": len(reader.pages)})

        elif ext in TEXT_EXT:
            with open(path, "r", encoding="utf-8", errors="ignore") as fh:
                text = fh.read(2 * 1024 * 1024)
            meta.update({"字符数": len(text)})

        elif ext in VIDEO_EXT:
            dur = read_mp4_duration(path)
            if dur is None:
                # 无法解析容器时，按 1.2 MB/分钟的经验码率粗估，仅用于分档
                dur = meta["大小MB"] / 1.2 * 60.0
                meta["时长来源"] = "按体积估算"
            else:
                meta["时长来源"] = "容器元数据"
            meta["时长秒"] = int(dur)
            text = os.path.splitext(name)[0]
            # 同名字幕/脚本文件一并纳入分析
            for sub_ext in (".srt", ".vtt", ".txt", ".md"):
                sub = os.path.splitext(path)[0] + sub_ext
                if os.path.isfile(sub):
                    with open(sub, "r", encoding="utf-8", errors="ignore") as fh:
                        text += "\n" + fh.read(1024 * 512)
                    meta["附带字幕"] = os.path.basename(sub)
                    break

        elif ext in IMAGE_EXT:
            text = os.path.splitext(name)[0]
            meta["说明"] = "图片仅按文件名与数量计量"

        else:
            # 未知类型：尝试按文本读，读不出就只用文件名
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as fh:
                    text = fh.read(512 * 1024)
            except Exception:
                text = ""
            if not text.strip():
                text = os.path.splitext(name)[0]
    except Exception as exc:                          # pragma: no cover
        meta["解析异常"] = str(exc)[:200]
        text = os.path.splitext(name)[0]

    return text, meta


# ==========================================================================
# 五、存储层（SQLite）
# ==========================================================================
SCHEMA = """
CREATE TABLE IF NOT EXISTS students (
    sid         TEXT PRIMARY KEY,
    name        TEXT,
    school      TEXT,
    college     TEXT,
    major       TEXT,
    class_name  TEXT,
    grade_raw   TEXT,
    grade_guess TEXT,
    updated_at  TEXT
);
CREATE TABLE IF NOT EXISTS aliases (
    alias TEXT PRIMARY KEY,      -- 规范化后的昵称或平台账号
    sid   TEXT NOT NULL,
    kind  TEXT
);
CREATE TABLE IF NOT EXISTS baselines (
    sid   TEXT NOT NULL,
    grade TEXT NOT NULL,
    dim   TEXT NOT NULL,
    score REAL NOT NULL,
    src   TEXT,
    PRIMARY KEY (sid, grade, dim)
);
CREATE TABLE IF NOT EXISTS evidences (
    fp        TEXT PRIMARY KEY,   -- 内容指纹，重复上传自动忽略
    sid       TEXT NOT NULL,
    grade     TEXT NOT NULL,
    channel   TEXT NOT NULL,
    dim       TEXT NOT NULL,
    amount    REAL NOT NULL,      -- 计量规模
    quality   REAL NOT NULL,      -- 质量系数 [0,1]
    relevance REAL NOT NULL,      -- 文化相关度 [0.3,1]
    label     TEXT,               -- 指标名称（展示用）
    detail    TEXT,               -- 证据摘要（展示用）
    batch_id  TEXT,
    created_at TEXT
);
CREATE INDEX IF NOT EXISTS idx_ev_sid ON evidences(sid, grade);
CREATE TABLE IF NOT EXISTS uploads (
    batch_id   TEXT PRIMARY KEY,
    kind       TEXT,
    filename   TEXT,
    sha256     TEXT,
    grade      TEXT,
    n_rows     INTEGER,
    n_new      INTEGER,
    n_dup      INTEGER,
    note       TEXT,
    created_at TEXT
);
CREATE INDEX IF NOT EXISTS idx_up_sha ON uploads(sha256);
"""


class Store:
    """线程安全的 SQLite 封装。"""

    def __init__(self, db_path: str):
        self.db_path = db_path
        os.makedirs(os.path.dirname(os.path.abspath(db_path)) or ".", exist_ok=True)
        self._local = threading.local()
        self._lock = threading.Lock()
        with self.conn() as cx:
            cx.executescript(SCHEMA)

    def conn(self) -> sqlite3.Connection:
        cx = getattr(self._local, "cx", None)
        if cx is None:
            cx = sqlite3.connect(self.db_path, timeout=30, check_same_thread=False)
            cx.row_factory = sqlite3.Row
            cx.execute("PRAGMA journal_mode=WAL")
            cx.execute("PRAGMA synchronous=NORMAL")
            self._local.cx = cx
        return cx

    # ---- 学生与别名 -----------------------------------------------------
    def upsert_student(self, sid: str, **fields: Any) -> None:
        if not sid:
            return
        cx = self.conn()
        row = cx.execute("SELECT sid FROM students WHERE sid=?", (sid,)).fetchone()
        cols = ("name", "school", "college", "major", "class_name",
                "grade_raw", "grade_guess")
        if row is None:
            vals = [fields.get(c) or "" for c in cols]
            cx.execute(
                f"INSERT INTO students(sid,{','.join(cols)},updated_at) "
                f"VALUES(?,{','.join('?' * len(cols))},?)",
                [sid] + vals + [now_iso()])
        else:
            sets, vals = [], []
            for c in cols:
                v = fields.get(c)
                if v:                                  # 只补空，不覆盖已有值
                    sets.append(f"{c}=COALESCE(NULLIF({c},''),?)")
                    vals.append(v)
            if sets:
                cx.execute(f"UPDATE students SET {','.join(sets)},updated_at=? "
                           f"WHERE sid=?", vals + [now_iso(), sid])
        cx.commit()

    def add_alias(self, alias: str, sid: str, kind: str = "nick") -> None:
        alias = norm_name(alias)
        if not alias or not sid:
            return
        cx = self.conn()
        cx.execute("INSERT OR IGNORE INTO aliases(alias,sid,kind) VALUES(?,?,?)",
                   (alias, sid, kind))
        cx.commit()

    def resolve(self, *candidates: Any) -> Optional[str]:
        """按 学号 -> 别名 的顺序解析出学生主键。"""
        cx = self.conn()
        for c in candidates:
            sid = norm_sid(c)
            if sid and cx.execute("SELECT 1 FROM students WHERE sid=?",
                                  (sid,)).fetchone():
                return sid
        for c in candidates:
            alias = norm_name(c)
            if not alias:
                continue
            row = cx.execute("SELECT sid FROM aliases WHERE alias=?",
                             (alias,)).fetchone()
            if row:
                return row["sid"]
        return None

    # ---- 证据 -----------------------------------------------------------
    def add_evidence(self, ev: "Evidence") -> bool:
        cx = self.conn()
        cur = cx.execute(
            "INSERT OR IGNORE INTO evidences"
            "(fp,sid,grade,channel,dim,amount,quality,relevance,label,detail,"
            " batch_id,created_at) VALUES(?,?,?,?,?,?,?,?,?,?,?,?)",
            (ev.fp, ev.sid, ev.grade, ev.channel, ev.dim, ev.amount, ev.quality,
             ev.relevance, ev.label, ev.detail, ev.batch_id, now_iso()))
        return cur.rowcount > 0

    def set_baseline(self, sid: str, grade: str, scores: Dict[str, float],
                     src: str = "问卷") -> None:
        cx = self.conn()
        for d, v in scores.items():
            if d in DIMS:
                cx.execute("INSERT OR REPLACE INTO baselines(sid,grade,dim,score,src)"
                           " VALUES(?,?,?,?,?)", (sid, grade, d, clamp(v), src))
        cx.commit()

    def commit(self) -> None:
        self.conn().commit()

    def log_upload(self, **kw: Any) -> None:
        cx = self.conn()
        cx.execute(
            "INSERT OR REPLACE INTO uploads"
            "(batch_id,kind,filename,sha256,grade,n_rows,n_new,n_dup,note,created_at)"
            " VALUES(?,?,?,?,?,?,?,?,?,?)",
            (kw.get("batch_id"), kw.get("kind"), kw.get("filename"),
             kw.get("sha256"), kw.get("grade"), kw.get("n_rows", 0),
             kw.get("n_new", 0), kw.get("n_dup", 0), kw.get("note", ""), now_iso()))
        cx.commit()


@dataclass
class Evidence:
    sid: str
    grade: str
    channel: str
    dim: str
    amount: float
    quality: float
    relevance: float
    label: str
    detail: str
    batch_id: str = ""
    fp: str = ""

    def finalize(self, key: str) -> "Evidence":
        """key 必须能唯一标识这条证据的"内容"，以保证重复上传幂等。"""
        self.fp = sha256_text(self.sid, self.grade, self.channel, self.dim, key)
        return self


# ==========================================================================
# 六、数据导入：自动识别 4 类导出表 + 问卷
# ==========================================================================
def _sheets(path: str) -> Dict[str, List[Tuple]]:
    """读出工作簿全部工作表为 {sheet_name: [row_tuple, ...]}。"""
    if openpyxl is None:
        raise RuntimeError("缺少 openpyxl，无法解析 Excel。请先 pip install openpyxl")
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    out: Dict[str, List[Tuple]] = {}
    for ws in wb.worksheets:
        out[ws.title] = list(ws.iter_rows(values_only=True))
    wb.close()
    return out


def _header_map(header_row: Sequence[Any]) -> Dict[str, int]:
    """列名 -> 列号（规范化后）。"""
    return {norm_text(h): i for i, h in enumerate(header_row) if h is not None}


def _pick(hm: Dict[str, int], *keys: str) -> Optional[int]:
    """按候选关键字模糊取列号。"""
    for k in keys:
        nk = norm_text(k)
        if nk in hm:
            return hm[nk]
    for k in keys:
        nk = norm_text(k)
        for name, idx in hm.items():
            if nk and nk in name:
                return idx
    return None


def detect_kind(path: str) -> str:
    """识别上传的 Excel 属于哪一类导出表。"""
    try:
        sh = _sheets(path)
    except Exception:
        return "unknown"
    names = " ".join(sh.keys())
    if "报名" in names:
        return "enrollment"
    if "作业" in names:
        return "homework"
    if "视频" in names:
        return "video"
    for rows in sh.values():
        if not rows:
            continue
        hm = _header_map(rows[0])
        joined = " ".join(hm.keys())
        if "用户提问" in joined or "ai回复" in joined:
            return "dialogue"
        if "报名课程" in joined or "文化课程" in joined:
            return "enrollment"
        if "需提交作业数" in joined:
            return "homework"
        if "视频总数" in joined:
            return "video"
        if "提交答卷时间" in joined or "答卷" in joined:
            return "questionnaire"
    return "unknown"


# ---- 6.1 报名表：建立学生主索引 -------------------------------------------
def import_enrollment(store: Store, path: str, grade_sel: Optional[str],
                      batch_id: str) -> Dict[str, Any]:
    sh = _sheets(path)
    n_new = n_dup = n_rows = 0
    unmatched: List[str] = []

    summary = None
    for name, rows in sh.items():
        if "汇总" in name or summary is None:
            summary = rows
            if "汇总" in name:
                break
    if not summary or len(summary) < 2:
        return {"n_rows": 0, "n_new": 0, "n_dup": 0, "note": "未找到有效数据行"}

    hm = _header_map(summary[0])
    c_nick = _pick(hm, "学生昵称", "姓名", "昵称")
    c_sid = _pick(hm, "学号")
    c_school = _pick(hm, "学校")
    c_col = _pick(hm, "学院")
    c_maj = _pick(hm, "专业")
    c_grade = _pick(hm, "年级")
    c_cls = _pick(hm, "行政班级", "班级")
    c_courses = _pick(hm, "报名课程", "文化课程")
    c_ncourse = _pick(hm, "报名课程数")

    def g(row: Sequence[Any], idx: Optional[int]) -> Any:
        return row[idx] if idx is not None and idx < len(row) else None

    for row in summary[1:]:
        if row is None or all(c is None for c in row):
            continue
        n_rows += 1
        nick = str(g(row, c_nick) or "").strip()
        sid = norm_sid(g(row, c_sid))
        if not sid:
            sid = "NICK:" + (norm_name(nick) or sha256_text(row)[:12])
            unmatched.append(nick)
        grade_raw = str(g(row, c_grade) or "")
        store.upsert_student(
            sid, name=nick, school=str(g(row, c_school) or ""),
            college=str(g(row, c_col) or ""), major=str(g(row, c_maj) or ""),
            class_name=str(g(row, c_cls) or ""), grade_raw=grade_raw,
            grade_guess=normalize_grade(grade_raw) or "")
        store.add_alias(nick, sid, "nick")

        grade = grade_sel or normalize_grade(grade_raw) or GRADES[0]
        raw_courses = str(g(row, c_courses) or "")
        courses = [c.strip() for c in re.split(r"[；;、\n]+", raw_courses) if c.strip()]
        if not courses and c_ncourse is not None:
            courses = [f"未命名课程{i + 1}" for i in range(to_int(g(row, c_ncourse)))]

        for course in courses[:12]:
            sig = analyze_text(course)
            rel = sig.relevance
            for dim in DIMS:
                w = CHANNEL_WEIGHTS["enrollment"][dim]
                if w <= 0:
                    continue
                ev = Evidence(sid=sid, grade=grade, channel="enrollment", dim=dim,
                              amount=2.0, quality=1.0, relevance=rel,
                              label="报名文化课程",
                              detail=f"《{course}》（文化相关度 {rel:.2f}）",
                              batch_id=batch_id).finalize(course)
                if store.add_evidence(ev):
                    n_new += 1
                else:
                    n_dup += 1
    store.commit()
    return {"n_rows": n_rows, "n_new": n_new, "n_dup": n_dup,
            "note": f"无学号学生 {len(unmatched)} 人（已按昵称建档）"}


# ---- 6.2 作业完成率 --------------------------------------------------------
def import_homework(store: Store, path: str, grade_sel: Optional[str],
                    batch_id: str) -> Dict[str, Any]:
    sh = _sheets(path)
    rows = None
    for name, r in sh.items():
        if "课程" in name and r:                 # 优先用带课程与学号的明细表
            rows = r
            break
    if rows is None:
        rows = next((r for r in sh.values() if r), None)
    if not rows or len(rows) < 2:
        return {"n_rows": 0, "n_new": 0, "n_dup": 0, "note": "未找到有效数据行"}

    hm = _header_map(rows[0])
    c_nick = _pick(hm, "学生昵称", "姓名")
    c_sid = _pick(hm, "学号")
    c_course = _pick(hm, "课程（班级）", "课程")
    c_need = _pick(hm, "需提交作业数")
    c_done = _pick(hm, "已提交数")
    c_rate = _pick(hm, "作业提交率")
    c_score = _pick(hm, "作业得分合计")
    c_full = _pick(hm, "作业满分合计")
    c_srate = _pick(hm, "得分率")

    n_new = n_dup = n_rows = n_skip = 0

    def g(row, idx):
        return row[idx] if idx is not None and idx < len(row) else None

    for row in rows[1:]:
        if row is None or all(c is None for c in row):
            continue
        n_rows += 1
        sid = store.resolve(g(row, c_sid), g(row, c_nick))
        if not sid:
            n_skip += 1
            continue
        done = to_int(g(row, c_done))
        if done <= 0:
            continue
        need = max(to_int(g(row, c_need)), done)
        submit_rate = parse_percent(g(row, c_rate)) or (done / need if need else 0.0)
        score_rate = parse_percent(g(row, c_srate))
        if not score_rate and c_full is not None:
            full = to_float(g(row, c_full))
            score_rate = to_float(g(row, c_score)) / full if full > 0 else 0.0
        course = str(g(row, c_course) or "课程作业")
        sig = analyze_text(course)
        rel = sig.relevance
        # 质量 = 提交率与得分率的加权；未打分的作业按提交行为给 0.6 保底
        quality = clamp(0.35 * submit_rate + 0.65 * (score_rate or 0.6))
        amount = min(float(done), 20.0) * 1.2 * (0.5 + 0.5 * (score_rate or 0.6))

        for dim in DIMS:
            if CHANNEL_WEIGHTS["homework"][dim] <= 0:
                continue
            ev = Evidence(sid=sid, grade=grade_sel or GRADES[0], channel="homework",
                          dim=dim, amount=amount, quality=quality, relevance=rel,
                          label="课程作业完成",
                          detail=(f"《{course}》提交 {done}/{need} 份，"
                                  f"得分率 {score_rate * 100:.0f}%"),
                          batch_id=batch_id).finalize(f"{course}|{done}|{need}|{score_rate:.4f}")
            if store.add_evidence(ev):
                n_new += 1
            else:
                n_dup += 1
    store.commit()
    return {"n_rows": n_rows, "n_new": n_new, "n_dup": n_dup,
            "note": f"未匹配到学生的记录 {n_skip} 条（建议先导入报名表）"}


# ---- 6.3 视频完成率 --------------------------------------------------------
def import_video(store: Store, path: str, grade_sel: Optional[str],
                 batch_id: str) -> Dict[str, Any]:
    sh = _sheets(path)
    rows = None
    for name, r in sh.items():
        if "课程" in name and r:
            rows = r
            break
    if rows is None:
        rows = next((r for r in sh.values() if r), None)
    if not rows or len(rows) < 2:
        return {"n_rows": 0, "n_new": 0, "n_dup": 0, "note": "未找到有效数据行"}

    hm = _header_map(rows[0])
    c_nick = _pick(hm, "学生昵称", "姓名")
    c_sid = _pick(hm, "学号")
    c_course = _pick(hm, "课程（班级）", "课程")
    c_total = _pick(hm, "视频总数")
    c_done = _pick(hm, "已完成视频数")
    c_rate = _pick(hm, "课程视频完成率", "总体视频完成率", "完成率")

    n_new = n_dup = n_rows = n_skip = 0

    def g(row, idx):
        return row[idx] if idx is not None and idx < len(row) else None

    for row in rows[1:]:
        if row is None or all(c is None for c in row):
            continue
        n_rows += 1
        sid = store.resolve(g(row, c_sid), g(row, c_nick))
        if not sid:
            n_skip += 1
            continue
        done = to_int(g(row, c_done))
        if done <= 0:
            continue
        total = max(to_int(g(row, c_total)), done)
        rate = parse_percent(g(row, c_rate)) or (done / total if total else 0.0)
        course = str(g(row, c_course) or "课程视频")
        rel = analyze_text(course).relevance
        amount = min(float(done), 40.0) / 4.0
        quality = clamp(0.4 + 0.6 * rate)

        for dim in DIMS:
            if CHANNEL_WEIGHTS["video"][dim] <= 0:
                continue
            ev = Evidence(sid=sid, grade=grade_sel or GRADES[0], channel="video",
                          dim=dim, amount=amount, quality=quality, relevance=rel,
                          label="课程视频学习",
                          detail=(f"《{course}》完成 {done}/{total} 个，"
                                  f"完成率 {rate * 100:.0f}%"),
                          batch_id=batch_id).finalize(f"{course}|{done}|{total}")
            if store.add_evidence(ev):
                n_new += 1
            else:
                n_dup += 1
    store.commit()
    return {"n_rows": n_rows, "n_new": n_new, "n_dup": n_dup,
            "note": f"未匹配到学生的记录 {n_skip} 条（建议先导入报名表）"}


# ---- 6.4 文化对话记录 ------------------------------------------------------
def import_dialogue(store: Store, path: str, grade_sel: Optional[str],
                    batch_id: str) -> Dict[str, Any]:
    sh = _sheets(path)
    rows = next((r for r in sh.values() if r), None)
    if not rows or len(rows) < 2:
        return {"n_rows": 0, "n_new": 0, "n_dup": 0, "note": "未找到有效数据行"}

    hm = _header_map(rows[0])
    c_nick = _pick(hm, "学生昵称", "姓名")
    c_acct = _pick(hm, "账号")
    c_time = _pick(hm, "提问时间")
    c_q = _pick(hm, "用户提问")
    c_a = _pick(hm, "ai回复", "AI回复", "回复")
    c_tok = _pick(hm, "消耗token", "token")

    n_new = n_dup = n_rows = n_skip = 0
    unmatched: set = set()

    def g(row, idx):
        return row[idx] if idx is not None and idx < len(row) else None

    for row in rows[1:]:
        if row is None or all(c is None for c in row):
            continue
        n_rows += 1
        nick = str(g(row, c_nick) or "")
        acct = str(g(row, c_acct) or "")
        sid = store.resolve(nick, acct)
        if not sid:
            unmatched.add(nick or acct)
            n_skip += 1
            continue
        if acct:
            store.add_alias(acct, sid, "account")

        question = str(g(row, c_q) or "")
        answer = str(g(row, c_a) or "")
        q_norm = norm_text(question)
        if not q_norm:
            continue
        # 口水轮次直接降权，不给"你好/继续"刷分空间
        filler = q_norm in LEX_FILLER or len(q_norm) < 4
        sig = analyze_text(question + "\n" + answer[:1500])
        rel = sig.relevance
        length_q = len(q_norm)
        quality = 0.12 if filler else clamp(0.35 + 0.65 * min(1.0, length_q / 60.0))
        tokens = to_float(g(row, c_tok))
        if tokens > 1500:
            quality = clamp(quality + 0.10)

        dist = dim_distribution(sig)
        stamp = str(g(row, c_time) or "")
        key = sha256_text(stamp, question[:400])

        for dim in DIMS:
            w = CHANNEL_WEIGHTS["dialogue"][dim]
            if w <= 0:
                continue
            amount = 0.45 * (0.4 + 0.6 * dist[dim] * len(DIMS) / 2.0)
            ev = Evidence(sid=sid, grade=grade_sel or GRADES[0], channel="dialogue",
                          dim=dim, amount=amount, quality=quality, relevance=rel,
                          label="与禾灵儿的文化问答",
                          detail=(f"{stamp[:16]}｜{question.strip()[:48]}"
                                  + ("…" if len(question.strip()) > 48 else "")),
                          batch_id=batch_id).finalize(key)
            if store.add_evidence(ev):
                n_new += 1
            else:
                n_dup += 1
    store.commit()
    return {"n_rows": n_rows, "n_new": n_new, "n_dup": n_dup,
            "note": f"未关联到学号的对话 {n_skip} 条，涉及 {len(unmatched)} 个昵称"}


# ---- 6.5 五维问卷 ----------------------------------------------------------
LIKERT_MAP = {
    "完全不符合": 0.0, "非常不符合": 0.0, "很不符合": 0.0, "完全不同意": 0.0,
    "不符合": 0.25, "比较不符合": 0.25, "不太符合": 0.25, "不同意": 0.25,
    "一般": 0.5, "中立": 0.5, "说不清": 0.5, "不确定": 0.5,
    "比较符合": 0.75, "基本符合": 0.75, "较符合": 0.75, "同意": 0.75,
    "完全符合": 1.0, "非常符合": 1.0, "很符合": 1.0, "完全同意": 1.0,
    "从不": 0.0, "很少": 0.25, "有时": 0.5, "经常": 0.75, "总是": 1.0,
}


def _likert_value(v: Any, scale_max: int = 5) -> Optional[float]:
    """把问卷作答转成 [0,1]。支持文字量表、1~5/1~7 数字、百分数。"""
    if v is None:
        return None
    if isinstance(v, (int, float)):
        x = float(v)
        if x <= 1.0:
            return clamp(x)
        return clamp((x - 1.0) / max(1.0, scale_max - 1.0))
    t = norm_text(v)
    if not t:
        return None
    for k, val in LIKERT_MAP.items():
        if norm_text(k) == t:
            return val
    for k, val in LIKERT_MAP.items():
        if norm_text(k) in t:
            return val
    m = re.search(r"(\d+(?:\.\d+)?)", t)
    if m:
        x = float(m.group(1))
        if "%" in t:
            return clamp(x / 100.0)
        if x <= 1.0:
            return clamp(x)
        return clamp((x - 1.0) / max(1.0, scale_max - 1.0))
    return None


def _question_dim(header: str, override: Dict[str, Dict[str, Any]]) -> Optional[Tuple[str, bool]]:
    """判定一道题属于哪个维度，返回 (维度, 是否反向计分)。

    优先级：外部映射文件 > 题干中直接出现维度名 > [E1]/C3 之类编码 > 词典语义归类。
    """
    raw = str(header or "")
    key = norm_text(raw)
    for pat, cfg in override.items():
        if norm_text(pat) and norm_text(pat) in key:
            dim = cfg.get("dim")
            if dim in DIMS:
                return dim, bool(cfg.get("reverse"))
    for dim in DIMS:
        if norm_text(dim) in key:
            return dim, False
    m = re.search(r"[\[\(（【]?\s*([ECIDR])\s*\d{1,2}\s*[\]\)）】]?", raw, re.I)
    if m:
        return CODE_DIM[m.group(1).upper()], False
    sig = analyze_text(raw)
    specificity = {"文化体验": 1.0, "文化认知": 1.0, "文化志趣": 0.8,
                   "文化认同": 1.0, "文化研创": 1.0}
    best, best_n = None, 0.0
    for dim, n in sig.dim_hits.items():
        v = n * specificity[dim]
        if v > best_n:
            best, best_n = dim, v
    if best_n > 0:
        return best, False
    return None


def import_questionnaire(store: Store, path: str, grade_sel: Optional[str],
                         batch_id: str) -> Dict[str, Any]:
    """解析问卷星等平台导出的 xlsx，生成五维初始得分。"""
    override: Dict[str, Dict[str, Any]] = {}
    cfg_path = os.path.join(config_dir(), "questionnaire_map.json")
    if os.path.isfile(cfg_path):
        try:
            with open(cfg_path, "r", encoding="utf-8") as fh:
                override = json.load(fh)
        except Exception:
            override = {}

    sh = _sheets(path)
    rows = next((r for r in sh.values() if r and len(r) > 1), None)
    if not rows:
        return {"n_rows": 0, "n_new": 0, "n_dup": 0, "note": "未找到有效数据行"}

    header = rows[0]
    hm = _header_map(header)
    c_sid = _pick(hm, "学号", "学生学号")
    c_name = _pick(hm, "姓名", "学生姓名", "昵称")
    c_cls = _pick(hm, "班级", "行政班级")

    # 若表内已有各维度总分列，直接采用
    direct: Dict[str, int] = {}
    for dim in DIMS:
        idx = _pick(hm, f"{dim}得分", f"{dim}分", dim)
        if idx is not None:
            direct[dim] = idx

    q_cols: Dict[int, Tuple[str, bool]] = {}
    if len(direct) < len(DIMS):
        skip = {"序号", "提交答卷时间", "所用时间", "来源", "来源详情", "来自ip",
                "总分", "答题时长", "用户id", "微信昵称"}
        for i, h in enumerate(header):
            if h is None or i in (c_sid, c_name, c_cls):
                continue
            if norm_text(h) in skip:
                continue
            res = _question_dim(str(h), override)
            if res:
                q_cols[i] = res

    n_students = 0
    dim_counts = {d: 0 for d in DIMS}
    for _, (d, _rev) in q_cols.items():
        dim_counts[d] += 1

    for row in rows[1:]:
        if row is None or all(c is None for c in row):
            continue
        sid_raw = row[c_sid] if c_sid is not None and c_sid < len(row) else None
        name_raw = row[c_name] if c_name is not None and c_name < len(row) else None
        sid = store.resolve(sid_raw, name_raw)
        if not sid:
            sid = norm_sid(sid_raw)
            if not sid:
                continue
            store.upsert_student(sid, name=str(name_raw or ""),
                                 class_name=str(row[c_cls]) if c_cls is not None
                                 and c_cls < len(row) and row[c_cls] else "")
            if name_raw:
                store.add_alias(str(name_raw), sid, "nick")

        scores: Dict[str, List[float]] = {d: [] for d in DIMS}
        if direct:
            for dim, idx in direct.items():
                if idx < len(row):
                    v = _likert_value(row[idx])
                    if v is not None:
                        scores[dim].append(v)
        for idx, (dim, rev) in q_cols.items():
            if idx >= len(row):
                continue
            v = _likert_value(row[idx])
            if v is None:
                continue
            scores[dim].append(1.0 - v if rev else v)

        final: Dict[str, float] = {}
        for dim in DIMS:
            vals = scores[dim]
            final[dim] = round(sum(vals) / len(vals), 4) if vals else BASELINE_DEFAULT
        if any(scores[d] for d in DIMS):
            store.set_baseline(sid, grade_sel or GRADES[0], final)
            n_students += 1
    store.commit()

    detail = "、".join(f"{d} {dim_counts[d]} 题" for d in DIMS) if q_cols else "使用维度得分列"
    return {"n_rows": len(rows) - 1, "n_new": n_students, "n_dup": 0,
            "note": f"识别到 {n_students} 名学生的问卷初始得分；题目归类：{detail}"}


# ---- 6.6 学生作业成果 / 自建智能体 -----------------------------------------
def import_artifact(store: Store, file_path: str, filename: str, sid: str,
                    grade: str, channel: str, batch_id: str,
                    note_text: str = "") -> Dict[str, Any]:
    """把一份学生作业或一个自建智能体转成证据。

    channel: 'artifact'（作业成果）或 'agent'（自建智能体）
    """
    with open(file_path, "rb") as fh:
        digest = sha256_bytes(fh.read())
    text, meta = extract_text_from_file(file_path, filename)
    full_text = (text or "") + "\n" + (note_text or "")
    sig = analyze_text(full_text, max_keywords=16)
    rel = sig.relevance
    ext = os.path.splitext(filename)[1].lower()

    # ---- 规模量：不同载体用各自的量纲折算成可比的"作品当量" ----
    if channel == "agent":
        base = 4.0
        base += min(2.0, len(full_text) / 1500.0)            # 提示词/知识库体量
    elif ext in VIDEO_EXT:
        secs = float(meta.get("时长秒") or 0)
        base = 2.0 + min(3.0, secs / 90.0)                   # 90 秒为一档
    elif ext in PPT_EXT:
        base = 1.5 + min(2.5, to_float(meta.get("页数")) / 8.0)
    elif ext in DOC_EXT or ext in PDF_EXT:
        base = 1.5 + min(2.5, len(full_text) / 1200.0)
    elif ext in XLS_EXT:
        base = 1.2 + min(2.0, to_float(meta.get("有效单元格")) / 400.0)
    elif ext in IMAGE_EXT:
        base = 1.0
    else:
        base = 1.2 + min(2.0, len(full_text) / 1500.0)

    # ---- 质量系数：文本充实度 + 结构完整度 ----
    quality = 0.45
    quality += min(0.25, len(full_text) / 4000.0)
    quality += min(0.15, (to_float(meta.get("插图数")) + to_float(meta.get("表格数"))) / 8.0)
    if sig.culture_hits >= 3:
        quality += 0.15
    quality = clamp(quality)

    dist = dim_distribution(sig)
    label = "自建文化智能体" if channel == "agent" else "学生作业成果"
    kw = "、".join(sig.keywords[:6]) or "未检出文化主题词"
    meta_bits = "，".join(f"{k} {v}" for k, v in meta.items()
                          if k not in ("文件名",) and v not in ("", None))
    detail = f"{filename}｜{meta_bits}｜主题词：{kw}"

    n_new = n_dup = 0
    for dim in DIMS:
        w = CHANNEL_WEIGHTS[channel][dim]
        if w <= 0:
            continue
        amount = base * (0.45 + 0.55 * dist[dim] * len(DIMS) / 2.0)
        ev = Evidence(sid=sid, grade=grade, channel=channel, dim=dim,
                      amount=amount, quality=quality, relevance=rel,
                      label=label, detail=detail, batch_id=batch_id
                      ).finalize(digest)          # 指纹只认文件内容，重传不叠加
        if store.add_evidence(ev):
            n_new += 1
        else:
            n_dup += 1
    store.commit()
    return {"n_rows": 1, "n_new": n_new, "n_dup": n_dup, "sha256": digest,
            "note": ("重复文件，未重复计分" if n_new == 0 else
                     f"已计入，文化相关度 {rel:.2f}，主题词：{kw}"),
            "meta": meta, "relevance": rel, "keywords": sig.keywords}


# ==========================================================================
# 七、打分引擎
# ==========================================================================
def _aggregate_A(store: Store, sid: str, grade: str) -> Tuple[Dict[str, float],
                                                              Dict[str, Dict[str, float]]]:
    """汇总某学生某年级的有效证据量 A_d，并返回按渠道拆分的明细。

    每个渠道对单一维度的贡献先各自截断到 CHANNEL_A_CAP，再求和。
    截断作用在"汇总值"而非"单条证据"上，因此与上传顺序完全无关。
    """
    cx = store.conn()
    rows = cx.execute(
        "SELECT channel, dim, SUM(amount*quality*relevance) AS s "
        "FROM evidences WHERE sid=? AND grade=? GROUP BY channel, dim",
        (sid, grade)).fetchall()
    per_channel: Dict[str, Dict[str, float]] = {}
    for r in rows:
        ch, dim = r["channel"], r["dim"]
        if dim not in DIMS:
            continue
        w = CHANNEL_WEIGHTS.get(ch, {}).get(dim, 0.0)
        val = float(r["s"] or 0.0) * w
        cap = CHANNEL_A_CAP.get(ch, 8.0)
        per_channel.setdefault(ch, {})[dim] = min(val, cap)
    A = {d: 0.0 for d in DIMS}
    for ch, m in per_channel.items():
        for d, v in m.items():
            A[d] += v
    return A, per_channel


def _explicit_baseline(store: Store, sid: str, grade: str) -> Optional[Dict[str, float]]:
    cx = store.conn()
    rows = cx.execute("SELECT dim, score FROM baselines WHERE sid=? AND grade=?",
                      (sid, grade)).fetchall()
    if not rows:
        return None
    b = {d: BASELINE_DEFAULT for d in DIMS}
    for r in rows:
        if r["dim"] in DIMS:
            b[r["dim"]] = float(r["score"])
    return b


def _has_evidence(store: Store, sid: str, grade: str) -> bool:
    cx = store.conn()
    r = cx.execute("SELECT 1 FROM evidences WHERE sid=? AND grade=? LIMIT 1",
                   (sid, grade)).fetchone()
    return r is not None


def compute_portrait(store: Store, sid: str) -> Dict[str, Any]:
    """计算一名学生四个年级的完整画像（含初始得分承接与成长轨迹）。"""
    grades: Dict[str, Any] = {}
    carry: Optional[Dict[str, float]] = None

    for g in GRADES:
        explicit = _explicit_baseline(store, sid, g)
        if explicit is not None:
            B, b_src = explicit, "本年级五维问卷"
        elif carry is not None:
            B, b_src = dict(carry), "承接上一学年画像"
        else:
            B, b_src = {d: BASELINE_DEFAULT for d in DIMS}, "默认初始得分 0.50"

        A, per_channel = _aggregate_A(store, sid, g)
        S: Dict[str, float] = {}
        rho: Dict[str, float] = {}
        for d in DIMS:
            r = 1.0 - math.exp(-A[d] / KAPPA[d]) if A[d] > 0 else 0.0
            rho[d] = r
            S[d] = round(clamp(B[d] + (1.0 - B[d]) * r), 4)

        has_data = explicit is not None or _has_evidence(store, sid, g)
        cci = round(sum(S[d] for d in DIMS) / len(DIMS), 4)
        grades[g] = {
            "grade": g, "has_data": has_data, "baseline": {d: round(B[d], 4) for d in DIMS},
            "baseline_source": b_src, "A": {d: round(A[d], 4) for d in DIMS},
            "rho": {d: round(rho[d], 4) for d in DIMS}, "scores": S, "cci": cci,
            "per_channel": {ch: {d: round(v, 4) for d, v in m.items()}
                            for ch, m in per_channel.items()},
        }
        if has_data:
            carry = S

    # 环比：与上一个"有数据"的年级比
    prev_cci = None
    for g in GRADES:
        item = grades[g]
        if prev_cci is None:
            item["delta"] = 0.0
            item["growth"] = 0.0
        else:
            item["delta"] = round(item["cci"] - prev_cci, 4)
            item["growth"] = round((item["cci"] - prev_cci) / prev_cci, 4) if prev_cci else 0.0
        if item["has_data"]:
            prev_cci = item["cci"]
    return grades


def student_row(store: Store, sid: str) -> Dict[str, Any]:
    cx = store.conn()
    r = cx.execute("SELECT * FROM students WHERE sid=?", (sid,)).fetchone()
    return dict(r) if r else {"sid": sid, "name": "", "class_name": ""}


# ---- 指标明细：详情页左上「提取到的指标信息」 -------------------------------
def indicator_lines(store: Store, sid: str, grade: str,
                    limit_per_channel: int = 4) -> List[Dict[str, Any]]:
    cx = store.conn()
    out: List[Dict[str, Any]] = []
    chans = cx.execute(
        "SELECT channel, COUNT(DISTINCT fp) AS n FROM evidences "
        "WHERE sid=? AND grade=? GROUP BY channel", (sid, grade)).fetchall()
    order = ["dialogue", "enrollment", "video", "homework", "artifact", "agent"]
    chans = sorted(chans, key=lambda r: order.index(r["channel"])
                   if r["channel"] in order else 99)
    for c in chans:
        ch = c["channel"]
        rows = cx.execute(
            "SELECT DISTINCT detail, label, MAX(amount*quality*relevance) AS s "
            "FROM evidences WHERE sid=? AND grade=? AND channel=? "
            "GROUP BY detail ORDER BY s DESC LIMIT ?",
            (sid, grade, ch, limit_per_channel)).fetchall()
        n_items = cx.execute(
            "SELECT COUNT(DISTINCT detail) AS n FROM evidences "
            "WHERE sid=? AND grade=? AND channel=?", (sid, grade, ch)).fetchone()["n"]
        out.append({
            "channel": ch,
            "label": CHANNEL_LABEL.get(ch, ch),
            "count": n_items,
            "items": [r["detail"] for r in rows],
            "more": max(0, n_items - len(rows)),
        })
    return out


# ---- 维度矩阵：详情页左下「指标 → 维度」 -----------------------------------
def dimension_matrix(store: Store, sid: str, grade: str,
                     portrait: Dict[str, Any]) -> List[Dict[str, Any]]:
    per_channel = portrait["per_channel"]
    A = portrait["A"]
    out = []
    for d in DIMS:
        parts = []
        for ch, m in per_channel.items():
            v = m.get(d, 0.0)
            if v > 0.01:
                parts.append((CHANNEL_LABEL.get(ch, ch), v))
        parts.sort(key=lambda x: -x[1])
        total = A[d] or 1.0
        out.append({
            "dim": d, "code": DIM_CODE[d], "score": portrait["scores"][d],
            "baseline": portrait["baseline"][d], "A": A[d], "rho": portrait["rho"][d],
            "sources": [{"name": n, "share": round(v / total, 3), "value": round(v, 3)}
                        for n, v in parts],
        })
    return out


# ---- 教师干预建议：详情页右下 ----------------------------------------------
SUGGESTION_BANK: Dict[str, Dict[str, str]] = {
    "文化体验": {
        "覆盖型": "该生几乎没有留下现场性、亲历性的学习痕迹。建议指派 1 次“运河文化带”或“南湖—月河历史街区”线下研学打卡，"
                    "并要求在禾灵儿上提交不少于 5 条文化意象记录；同步开放 1 条虚拟漫游线路作为线上补充。",
        "效能型": "该生有参与但停留较浅（视频完成率或任务完成度偏低）。建议把大课任务拆成 3 次小任务分周推进，"
                    "每次配 1 个必答的现场观察问题，把“到过”变成“看懂”。",
    },
    "文化认知": {
        "覆盖型": "缺少可评的知识性证据。建议布置 1 次“嘉兴文脉关键词”闯关测验，并要求就 1 个薄弱知识点与禾灵儿完成 3 轮追问式问答后提交小结。",
        "效能型": "知识点掌握零散、迁移不足。建议采用“一图一史”作业：让学生用一张概念图串起江南文化的时间线与地理线，教师批注后二次修改。",
    },
    "文化志趣": {
        "覆盖型": "自主学习行为很少，基本是被动完成任务。建议设置“7 天文化打卡”轻量任务（每天在资讯中心读 1 篇并收藏），先把习惯建立起来。",
        "效能型": "有兴趣但不持久、主题过窄。建议用兴趣锚点扩面：从该生已关注的主题出发，推荐 2 个相邻主题的智能体，鼓励跨主题提问。",
    },
    "文化认同": {
        "覆盖型": "尚未观察到明确的情感归属与价值表达。建议安排 1 次“我身边的红船精神”口述作业，要求结合本人或家人的真实经历，并参与同伴互评。",
        "效能型": "表达停留在结论层面，缺少论证与个人立场。建议开展 1 次小型辩论或“向外国朋友介绍嘉兴”的双语短讲，倒逼其组织论据。",
    },
    "文化研创": {
        "覆盖型": "没有可评的创作或研究成果。建议定制 1 次低门槛微创作：用禾灵儿视频生成完成 1 部 60 秒文化微视频，5 天内提交并发布到视频广场。",
        "效能型": "有作品但原创度或规范性不够。建议开展 1 次“跨文化叙事再创作”：要求提供原创脚本、标注素材出处，并在 4 天内提交创意方案与修改说明。",
    },
}


def teacher_advice(portrait: Dict[str, Any], stu: Dict[str, Any]) -> List[str]:
    """依据五维得分与归因类型生成给授课教师的改进建议。"""
    S, A, rho = portrait["scores"], portrait["A"], portrait["rho"]
    ordered = sorted(DIMS, key=lambda d: S[d])
    advice: List[str] = []

    for d in ordered[:2]:
        # 归因分两类：覆盖型＝过程证据太少，样本不足以判断；
        #             效能型＝证据量已够，问题出在达成水平。
        cause = "覆盖型" if A[d] < 1.2 else "效能型"
        pct = rho[d] * 100
        why = (f"过程证据覆盖度仅 {pct:.0f}%，现有样本量尚不足以支撑稳定判断，属采集侧受限"
               if cause == "覆盖型" else
               f"过程证据覆盖度 {pct:.0f}%，样本量已过判定阈值，短板源于达成水平而非参与量")
        head = f"【{d} {S[d]:.2f} · 优先干预项】归因：{cause}短板 —— {why}。"
        advice.append(head + SUGGESTION_BANK[d][cause])

    top = ordered[-1]
    advice.append(
        f"【{top} {S[top]:.2f} · 相对优势项】该生在{top}上表现突出，建议赋予拓展性角色——"
        f"担任小组{('研学领队' if top == '文化体验' else '创作主创' if top == '文化研创' else '朋辈领学')}，"
        f"用其长板带动小组其余同学的短板。")

    spread = max(S.values()) - min(S.values())
    if spread >= 0.25:
        advice.append(
            f"【整体研判】五维极差 {spread:.2f}，结构性失衡明显。建议本学期只抓"
            f"{ordered[0]}一个维度，集中投放任务与资源，避免样样都补、样样不深。")
    elif portrait["cci"] < 0.60:
        advice.append("【整体研判】综合指数偏低且各维度普遍偏弱，属整体参与总量不足而非单项薄弱。"
                      "建议先解决“上平台”的问题：约定每周固定 1 次禾灵儿学习时段，先把行为量提上来。")
    else:
        advice.append("【整体研判】五维发展结构较为均衡，建议保持现有节奏，"
                      "下阶段可尝试把单项优势转化为公开成果（作品发布、朋辈分享）。")
    return advice


ALL_CHANNELS = ("dialogue", "enrollment", "video", "homework", "artifact", "agent")
COMPLETENESS_W = {"questionnaire": 0.25, "channels": 0.30, "rho": 0.25, "volume": 0.20}


def completeness_of(store: Store, sid: str,
                    grades: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """材料完整度：这份画像有多少证据在撑着。

    完整度低不等于分低，而是"证据不足、结论不可靠"，因此列表默认把材料最全的
    学生排在最前面，教师可以先看结论最可信的那批人。

        完整度 = 0.25·问卷 + 0.30·渠道覆盖 + 0.25·证据充分度 + 0.20·证据体量
    """
    grades = grades or compute_portrait(store, sid)
    cx = store.conn()
    # 注意：evidences 表里一条逻辑证据会展开成 5 行（每个维度一行），
    # 这里必须按 detail 去重，才能与网页版"证据条数"的口径一致。
    rows = cx.execute(
        "SELECT channel, COUNT(DISTINCT detail) AS n FROM evidences "
        "WHERE sid=? GROUP BY channel", (sid,)).fetchall()
    chans = {r["channel"] for r in rows if r["n"]}
    items = sum(r["n"] for r in rows)
    has_q = cx.execute("SELECT 1 FROM baselines WHERE sid=? LIMIT 1",
                       (sid,)).fetchone() is not None
    best_rho = 0.0
    for g in GRADES:
        pr = grades[g]
        if pr["has_data"]:
            m = sum(pr["rho"][d] for d in DIMS) / len(DIMS)
            best_rho = max(best_rho, m)
    cov = len(chans) / len(ALL_CHANNELS)
    vol = min(1.0, math.log(1 + items) / math.log(1 + 60))
    score = (COMPLETENESS_W["questionnaire"] * (1.0 if has_q else 0.0)
             + COMPLETENESS_W["channels"] * cov
             + COMPLETENESS_W["rho"] * best_rho
             + COMPLETENESS_W["volume"] * vol)
    missing = ([] if has_q else ["问卷初始得分"]) + \
              [CHANNEL_LABEL[c] for c in ALL_CHANNELS if c not in chans]
    return {"score": round(score, 4), "has_q": has_q, "channels": len(chans),
            "items": items, "rho": round(best_rho, 4), "missing": missing}


def completeness_tag(score: float) -> Tuple[str, str]:
    # 用界面上显示的整数百分比判档，避免出现“显示 50% 却标偏少”的割裂感
    score = round(score, 2)
    if score >= 0.75:
        return "材料完整", "cp-a"
    if score >= 0.50:
        return "较完整", "cp-b"
    if score >= 0.28:
        return "偏少", "cp-c"
    return "严重不足", "cp-d"


def level_of(x: float) -> str:
    if x >= 0.85:
        return "A 引领级"
    if x >= 0.70:
        return "B 胜任级"
    if x >= 0.55:
        return "C 发展级"
    return "D 起步级"


# ==========================================================================
# 八、图形（服务端生成内联 SVG，无外部依赖、离线可用）
# ==========================================================================
def radar_svg(scores: Dict[str, float], baseline: Dict[str, float],
              size: int = 340) -> str:
    """五维雷达图。实线为当前得分，虚线为初始得分。

    画布刻意做成横向更宽的矩形：左右两侧要放下 4 个汉字的维度名与分值，
    正方形画布会把标签裁掉。
    """
    w, h = int(size * 1.50), int(size * 1.20)
    cx, cy = w / 2.0, h / 2.0 + 4
    R = size * 0.34
    n = len(DIMS)
    pts_s, pts_b, axes, labels, dots = [], [], [], [], []
    for i, d in enumerate(DIMS):
        ang = -math.pi / 2 + 2 * math.pi * i / n
        ux, uy = math.cos(ang), math.sin(ang)
        axes.append(f'<line x1="{cx:.1f}" y1="{cy:.1f}" '
                    f'x2="{cx + ux * R:.1f}" y2="{cy + uy * R:.1f}" '
                    f'stroke="#e2e8f0" stroke-width="1"/>')
        sv = clamp(scores.get(d, 0.0))
        bv = clamp(baseline.get(d, 0.0))
        sx, sy = cx + ux * R * sv, cy + uy * R * sv
        pts_s.append(f"{sx:.1f},{sy:.1f}")
        pts_b.append(f"{cx + ux * R * bv:.1f},{cy + uy * R * bv:.1f}")
        dots.append(f'<circle cx="{sx:.1f}" cy="{sy:.1f}" r="4" fill="#059669"/>')
        lx, ly = cx + ux * (R + 40), cy + uy * (R + 27)
        if ux > 0.25:
            anchor = "start"
        elif ux < -0.25:
            anchor = "end"
        else:
            anchor = "middle"
        labels.append(
            f'<text x="{lx:.1f}" y="{ly:.1f}" text-anchor="{anchor}" '
            f'font-size="16" font-weight="700" fill="#1f2937">{d}</text>'
            f'<text x="{lx:.1f}" y="{ly + 21:.1f}" text-anchor="{anchor}" '
            f'font-size="19" font-weight="800" fill="#0f766e">{sv:.2f}</text>')
    rings = "".join(
        '<polygon points="{}" fill="none" stroke="#eef2f6" stroke-width="1"/>'.format(
            " ".join(
                f"{cx + math.cos(-math.pi / 2 + 2 * math.pi * i / n) * R * k:.1f},"
                f"{cy + math.sin(-math.pi / 2 + 2 * math.pi * i / n) * R * k:.1f}"
                for i in range(n)))
        for k in (0.25, 0.5, 0.75, 1.0))
    return (
        f'<svg viewBox="0 0 {w} {h}" width="100%" style="max-width:{w}px" '
        f'xmlns="http://www.w3.org/2000/svg">{rings}{"".join(axes)}'
        f'<polygon points="{" ".join(pts_b)}" fill="none" stroke="#94a3b8" '
        f'stroke-width="1.4" stroke-dasharray="5 4"/>'
        f'<polygon points="{" ".join(pts_s)}" fill="#10b981" fill-opacity="0.20" '
        f'stroke="#059669" stroke-width="2.6"/>'
        + "".join(dots) + "".join(labels) + "</svg>")


def trajectory_svg(grades: Dict[str, Any], width: int = 640, height: int = 260) -> str:
    """四学年成长轨迹折线图。"""
    pad_l, pad_r, pad_t, pad_b = 46, 16, 18, 34
    pw, ph = width - pad_l - pad_r, height - pad_t - pad_b
    colors = {"文化体验": "#0f766e", "文化认知": "#10b981", "文化志趣": "#0ea5e9",
              "文化认同": "#f59e0b", "文化研创": "#ef4444"}
    xs = [pad_l + pw * i / max(1, len(GRADES) - 1) for i in range(len(GRADES))]

    def y_of(v: float) -> float:
        return pad_t + ph * (1.0 - clamp(v))

    parts = [f'<svg viewBox="0 0 {width} {height}" width="100%" '
             f'style="max-width:{width}px;display:block;margin:0 auto" '
             f'xmlns="http://www.w3.org/2000/svg">']
    for k in (0.0, 0.25, 0.5, 0.75, 1.0):
        y = y_of(k)
        parts.append(f'<line x1="{pad_l}" y1="{y:.1f}" x2="{width - pad_r}" '
                     f'y2="{y:.1f}" stroke="#eef2f6" stroke-width="1"/>')
        parts.append(f'<text x="{pad_l - 8}" y="{y + 4:.1f}" text-anchor="end" '
                     f'font-size="10" fill="#94a3b8">{k:.2f}</text>')
    for i, g in enumerate(GRADES):
        parts.append(f'<text x="{xs[i]:.1f}" y="{height - 12}" text-anchor="middle" '
                     f'font-size="12" fill="#475569">{g}</text>')

    active = [i for i, g in enumerate(GRADES) if grades[g]["has_data"]]
    if not active:
        parts.append(f'<text x="{width / 2}" y="{height / 2}" text-anchor="middle" '
                     f'font-size="13" fill="#94a3b8">暂无可用于绘制轨迹的数据</text></svg>')
        return "".join(parts)

    for d in DIMS:
        pts = [f"{xs[i]:.1f},{y_of(grades[GRADES[i]]['scores'][d]):.1f}" for i in active]
        parts.append(f'<polyline points="{" ".join(pts)}" fill="none" '
                     f'stroke="{colors[d]}" stroke-width="1.8"/>')
        for p in pts:
            x, y = p.split(",")
            parts.append(f'<circle cx="{x}" cy="{y}" r="2.8" fill="{colors[d]}"/>')
    pts = [f"{xs[i]:.1f},{y_of(grades[GRADES[i]]['cci']):.1f}" for i in active]
    parts.append(f'<polyline points="{" ".join(pts)}" fill="none" stroke="#1e293b" '
                 f'stroke-width="2.6" stroke-dasharray="6 3"/>')

    lx = pad_l
    for d in DIMS:
        parts.append(f'<rect x="{lx}" y="{pad_t - 12}" width="10" height="3" '
                     f'fill="{colors[d]}"/>'
                     f'<text x="{lx + 14}" y="{pad_t - 7}" font-size="10" '
                     f'fill="#475569">{d}</text>')
        lx += 78
    parts.append(f'<rect x="{lx}" y="{pad_t - 12}" width="10" height="3" fill="#1e293b"/>'
                 f'<text x="{lx + 14}" y="{pad_t - 7}" font-size="10" '
                 f'fill="#475569">综合指数</text>')
    parts.append("</svg>")
    return "".join(parts)


# ==========================================================================
# 九、学习报告（Word）
# ==========================================================================
def build_report_docx(store: Store, sid: str, grade: str) -> bytes:
    """生成该生该学年的文化德育学习报告（.docx 字节流）。"""
    if _docx is None:
        raise RuntimeError('缺少 python-docx，无法导出报告。请先 pip install python-docx')
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Cm, Pt

    stu = student_row(store, sid)
    all_g = compute_portrait(store, sid)
    pr = all_g[grade]
    matrix = dimension_matrix(store, sid, grade, pr)
    lines = indicator_lines(store, sid, grade, limit_per_channel=6)
    advice = teacher_advice(pr, stu)

    doc = Document()
    sec = doc.sections[0]
    sec.page_width, sec.page_height = Cm(21.0), Cm(29.7)
    sec.top_margin = sec.bottom_margin = Cm(2.54)
    sec.left_margin = sec.right_margin = Cm(3.17)

    def run_style(run, cn='宋体', size=12.0, bold=False):
        run.font.size = Pt(size)
        run.bold = bold
        rPr = run._element.get_or_add_rPr()
        rf = rPr.find(qn('w:rFonts'))
        if rf is None:
            rf = OxmlElement('w:rFonts')
            rPr.insert(0, rf)
        en = 'Times New Roman' if cn == '宋体' else cn
        for attr in ('w:ascii', 'w:hAnsi', 'w:cs'):
            rf.set(qn(attr), en)
        rf.set(qn('w:eastAsia'), cn)

    def para(text='', cn='宋体', size=12.0, bold=False, align=None, indent=False):
        p = doc.add_paragraph()
        pf = p.paragraph_format
        pf.space_before = Pt(0)
        pf.space_after = Pt(0)
        pf.line_spacing = 1.5
        if align is not None:
            p.alignment = align
        if indent:
            pf.first_line_indent = Pt(size * 2)
        if text:
            run_style(p.add_run(text), cn, size, bold)
        return p

    def heading(text, size=14.0):
        p = para(text, cn='黑体', size=size, bold=True)
        p.paragraph_format.space_before = Pt(10)
        p.paragraph_format.space_after = Pt(4)
        return p

    para(f'{stu.get("name") or "该生"}文化德育数字画像学习报告', cn='黑体', size=16.0,
         bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    para(f'{grade}学年　｜　生成时间 {now_iso()}', cn='楷体', size=10.5,
         align=WD_ALIGN_PARAGRAPH.CENTER)
    para()

    heading('一、基本信息')
    info = (f'学号：{sid}　　姓名：{stu.get("name") or "—"}　　'
            f'班级：{stu.get("class_name") or "—"}\n'
            f'学院：{stu.get("college") or "—"}　　专业：{stu.get("major") or "—"}')
    for ln in info.split('\n'):
        para(ln, indent=True)

    heading('二、五维素养得分')
    tbl = doc.add_table(rows=1, cols=5)
    tbl.style = 'Table Grid'
    for j, h in enumerate(['维度', '初始得分', '本学年得分', '增幅', '等级']):
        c = tbl.rows[0].cells[j]
        c.text = ''
        run_style(c.paragraphs[0].add_run(h), '黑体', 10.5, False)
        c.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    for d in DIMS:
        row = tbl.add_row()
        vals = [d, f'{pr["baseline"][d]:.2f}', f'{pr["scores"][d]:.2f}',
                f'+{pr["scores"][d] - pr["baseline"][d]:.2f}', level_of(pr['scores'][d])]
        for j, v in enumerate(vals):
            c = row.cells[j]
            c.text = ''
            run_style(c.paragraphs[0].add_run(v), '宋体', 10.5, False)
            c.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    para()
    para(f'综合文化素养指数 CCI = {pr["cci"]:.2f}（{level_of(pr["cci"])}），'
         f'较上一学年 {pr["delta"]:+.2f}。', indent=True)

    heading('三、支撑得分的过程证据')
    if not lines:
        para('本学年暂未采集到过程性证据。', indent=True)
    for ln in lines:
        para(f'{ln["label"]}（{ln["count"]} 项）', cn='黑体', size=12.0)
        for item in ln['items']:
            para(f'　·　{item}', size=10.5)
        if ln['more']:
            para(f'　…… 另有 {ln["more"]} 项未列出', size=10.5)

    heading('四、维度归因')
    for m in matrix:
        src = '；'.join(f'{s["name"]} {s["share"] * 100:.0f}%' for s in m['sources'][:4])
        para(f'{m["dim"]}：得分 {m["score"]:.2f}，证据充分度 {m["rho"] * 100:.0f}%。'
             f'主要证据来源：{src or "暂无"}。', indent=True)

    heading('五、给授课教师的改进建议')
    for a in advice:
        para(a, indent=True)

    heading('六、说明')
    para('本报告由禾灵儿文化德育数字画像系统自动生成。五个维度取值区间为 0 至 1，'
         '初始得分来自开课第一节课的五维问卷，其后依据学生在平台上的课程学习、作业提交、'
         '文化对话、作品创作与智能体开发等过程证据动态累加，得分随材料增加而单调上升'
         '并以 1.00 为上限。全部计分规则公开、可复算，每一项得分均可追溯到具体的学习事件。',
         indent=True)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ==========================================================================
# 十、Web 应用（FastAPI）
# ==========================================================================
def default_db_path() -> str:
    env = os.environ.get('HLR_DB')
    if env:
        return env
    return os.path.join(config_dir(), 'hlr_data', 'portrait.db')


STORE = Store(default_db_path())

try:
    from fastapi import FastAPI, File, Form, HTTPException, Query, UploadFile
    from fastapi.responses import HTMLResponse, JSONResponse, Response
except Exception:                                      # pragma: no cover
    FastAPI = None

IMPORTERS = {
    'questionnaire': import_questionnaire,
    'enrollment': import_enrollment,
    'homework': import_homework,
    'video': import_video,
    'dialogue': import_dialogue,
}
KIND_LABEL = {
    'questionnaire': '五维问卷', 'enrollment': '学生报名情况',
    'homework': '学生作业完成率', 'video': '学生视频完成率',
    'dialogue': '文化对话记录',
}


# --------------------------------------------------------------------------
# 10.1 页面样式与骨架
# --------------------------------------------------------------------------
CSS = r'''
*{box-sizing:border-box;margin:0;padding:0}
body{font-family:"Microsoft YaHei","PingFang SC","Noto Sans CJK SC",system-ui,sans-serif;
     background:#f2f7f5;color:#1f2937;font-size:15.5px;line-height:1.7}
a{color:inherit;text-decoration:none}
.wrap{max-width:1280px;margin:0 auto;padding:20px}
.top{background:linear-gradient(90deg,#0f766e,#10b981);color:#fff;padding:18px 0}
.top .wrap{display:flex;align-items:center;justify-content:space-between;padding:0 20px}
.top h1{font-size:19px;font-weight:700}
.top .sub{font-size:12px;opacity:.9;margin-top:2px}
.top .rt{font-size:12px;opacity:.92;text-align:right}
.card{background:#fff;border-radius:12px;box-shadow:0 1px 3px rgba(15,23,42,.07);
      padding:18px 20px;margin-bottom:16px}
.card h2{font-size:17px;font-weight:800;margin-bottom:12px;display:flex;
         align-items:center;gap:8px;color:#0f766e}
.card h2 .n{background:#10b981;color:#fff;width:20px;height:20px;border-radius:5px;
            display:inline-flex;align-items:center;justify-content:center;font-size:12px}
table{width:100%;border-collapse:collapse}
th{background:#f8fafc;color:#334155;font-weight:700;font-size:14.5px;
   padding:12px 10px;text-align:center;border-bottom:1px solid #e2e8f0}
td{padding:11px 10px;text-align:center;border-bottom:1px solid #f1f5f9;
   font-size:14.5px;font-weight:600}
tbody tr:hover{background:#f0fdf9}
.btn{display:inline-block;background:#10b981;color:#fff;border:none;border-radius:999px;
     padding:8px 19px;font-size:14px;font-weight:600;cursor:pointer;transition:.15s}
.btn:hover{background:#0d9488}
.btn.ghost{background:#fff;color:#0f766e;border:1px solid #99f6e4}
.btn.ghost:hover{background:#f0fdfa}
.btn[disabled]{background:#e2e8f0;color:#94a3b8;cursor:not-allowed}
.btn.sm{padding:4px 14px;font-size:12px}
input,select,textarea{font-family:inherit;font-size:14px;padding:8px 12px;
     border:1px solid #d7e3de;border-radius:8px;background:#fff;color:#1f2937}
input:focus,select:focus,textarea:focus{outline:2px solid #99f6e4;outline-offset:-1px}
.row{display:flex;gap:10px;flex-wrap:wrap;align-items:center}
.grid2{display:grid;grid-template-columns:1fr 1fr;gap:16px;align-items:stretch}
/* 同一行的两块卡片等高，底边对齐；卡片内用 flex 让主体内容撑满剩余高度 */
.grid2>.card{display:flex;flex-direction:column}
.rwrap{flex:1;display:flex;gap:14px;flex-wrap:wrap;
       align-items:center;justify-content:space-between}
@media(max-width:940px){.grid2{grid-template-columns:1fr}}
.muted{color:#5b6b7c;font-size:13px}
.pill{display:inline-block;padding:3px 10px;border-radius:6px;font-size:13px;font-weight:800}
.p-E{background:#fff7ed;color:#c2410c}.p-C{background:#eff6ff;color:#1d4ed8}
.p-I{background:#f0fdf4;color:#15803d}.p-D{background:#fef2f2;color:#b91c1c}
.p-R{background:#f5f3ff;color:#6d28d9}
.ev{border-left:4px solid #6ee7b7;padding:6px 0 6px 12px;margin-bottom:11px;
     line-height:1.6}
.ev b{color:#0f766e;font-size:16.5px;font-weight:800;letter-spacing:.3px}
.ev .it{color:#243b53;font-size:15px;font-weight:600;display:block;margin-top:3px}
.sum{font-size:15px;font-weight:600;color:#334155;margin-bottom:11px;
     background:#f0fdfa;border-radius:8px;padding:8px 12px}
.sum b{color:#0f766e;font-size:17px;font-weight:800}
.radarnote{font-size:14px;font-weight:600;color:#475569;margin-top:8px}
.cci{background:linear-gradient(160deg,#ecfdf5,#d1fae5);border-radius:12px;padding:18px 20px;
     text-align:center;min-width:172px}
.cci .v{font-size:46px;font-weight:900;color:#065f46;line-height:1.05}
.adv{background:#f0fdf4;border-radius:9px;padding:12px 14px;margin-bottom:10px;
     font-size:14px;line-height:1.72;border-left:4px solid #34d399;color:#14342b}
.up{border:1.5px dashed #99f6e4;border-radius:10px;padding:14px;background:#f7fefc}
/* 解析日志默认收起，只在出错或用户主动展开时显示 */
.log{max-height:210px;overflow:auto;background:#0f172a;color:#c9f7e5;border-radius:8px;
     padding:10px 12px;font-size:12px;font-family:ui-monospace,Consolas,monospace;
     white-space:pre-wrap;line-height:1.6;margin-top:10px}
.statusbar{display:none;align-items:center;gap:10px;margin-top:10px;
           font-size:14px;font-weight:700;color:#0f766e}
.statusbar.on{display:flex}
.statusbar.err{color:#b91c1c}
.statusbar .dot{width:8px;height:8px;border-radius:50%;background:#10b981;flex:none}
.statusbar.err .dot{background:#ef4444}
.statusbar .lnk{margin-left:auto;color:#94a3b8;font-size:12.5px;font-weight:400;
                text-decoration:underline;cursor:pointer}
.statusbar .lnk:hover{color:#0f766e}
.tabs{display:flex;gap:6px;margin-bottom:14px;flex-wrap:wrap}
.tab{padding:7px 16px;border-radius:999px;background:#fff;border:1px solid #d7e3de;
     cursor:pointer;font-size:13px}
.tab.on{background:#0f766e;color:#fff;border-color:#0f766e}
.bar{height:6px;background:#eef2f6;border-radius:3px;overflow:hidden;margin-top:4px}
.bar>i{display:block;height:100%;background:linear-gradient(90deg,#34d399,#0f766e)}
.nodata{color:#94a3b8;font-size:13px}
.hide{display:none!important}
.cp{display:flex;flex-direction:column;gap:3px;align-items:stretch;text-align:left}
.cp .hd{display:flex;justify-content:space-between;align-items:baseline;gap:6px}
.cp .pc{font-weight:800;font-size:15px;color:#0f766e}
.cp .tg{font-size:12px;font-weight:700;padding:1px 8px;border-radius:5px}
.cp-a{background:#dcfce7;color:#15803d}.cp-b{background:#e0f2fe;color:#0369a1}
.cp-c{background:#fef3c7;color:#b45309}.cp-d{background:#fee2e2;color:#b91c1c}
.cp .ms{font-size:11.5px;color:#94a3b8;font-weight:400;line-height:1.4}
'''


def page(title: str, body_html: str, extra_js: str = '') -> str:
    return f'''<!doctype html><html lang="zh-CN"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{title} · {APP_NAME}</title><style>{CSS}</style></head><body>
<div class="top"><div class="wrap"><div>
  <h1>禾灵儿 · 文化德育数字画像</h1>
  <div class="sub">五维素养动态评价 ｜ 文化体验 · 文化认知 · 文化志趣 · 文化认同 · 文化研创</div>
</div><div class="rt">嘉兴大学 &nbsp;｜&nbsp; v{APP_VERSION}</div></div></div>
<div class="wrap">{body_html}</div>
<script>{extra_js}</script></body></html>'''


# --------------------------------------------------------------------------
# 10.2 路由
# --------------------------------------------------------------------------
if FastAPI is not None:
    app = FastAPI(title=APP_NAME, version=APP_VERSION)

    # ---------------- 列表页 ----------------
    @app.get('/', response_class=HTMLResponse)
    def index() -> str:
        body = r'''
<div class="card">
  <h2><span class="n">1</span>数据上传</h2>
  <div class="up">
    <div class="row" style="margin-bottom:10px">
      <label>材料类别</label>
      <select id="kind">
        <option value="auto">自动识别（推荐）</option>
        <option value="questionnaire">五维问卷</option>
        <option value="enrollment">学生报名情况</option>
        <option value="homework">学生作业完成率</option>
        <option value="video">学生视频完成率</option>
        <option value="dialogue">文化对话记录</option>
      </select>
      <label>年级</label>
      <select id="grade"><option>大一</option><option>大二</option>
        <option>大三</option><option>大四</option></select>
      <input type="file" id="files" multiple accept=".xlsx,.xlsm,.xls">
      <button class="btn" onclick="up()">上传并解析</button>
    </div>
    <div class="muted">首次使用请先上传<b>学生报名情况</b>建立学号索引，再上传其余表格；
      问卷用于生成五维初始得分。重复上传同一文件不会重复计分。</div>
  </div>
  <div class="statusbar" id="status">
    <span class="dot"></span><span id="statusText"></span>
    <span class="lnk" id="toggleLog" onclick="toggleLog()">查看详情</span>
  </div>
  <div class="log hide" id="log"></div>
</div>

<div class="card">
  <h2><span class="n">2</span>学生画像总览</h2>
  <div class="row" style="margin-bottom:12px">
    <input id="q" placeholder="搜索 班级 / 学号 / 姓名" style="width:250px"
           onkeydown="if(event.key==='Enter')load(1)">
    <button class="btn" onclick="load(1)">查询</button>
    <label class="muted">排序</label>
    <select id="sortBy" onchange="load(1)">
      <option value="complete">默认</option>
      <option value="cci">综合指数从高到低</option>
      <option value="cci_asc">综合指数从低到高</option>
      <option value="class">学号</option>
    </select>
    <span class="muted" id="cnt"></span>
  </div>
  <table><thead><tr><th>学号</th><th>姓名</th>
    <th>大一</th><th>大二</th><th>大三</th><th>大四</th></tr></thead>
    <tbody id="tb"><tr><td colspan="6" class="nodata">尚无数据，请先上传材料</td></tr></tbody>
  </table>
  <div class="row" style="margin-top:12px;justify-content:center">
    <button class="btn ghost sm" onclick="load(page-1)">上一页</button>
    <span class="muted" id="pg"></span>
    <button class="btn ghost sm" onclick="load(page+1)">下一页</button>
  </div>
</div>'''
        js = r'''
var page=1;
// 解析明细只写进隐藏的日志区；界面上只呈现一行状态
function log(t){var l=document.getElementById('log');
  l.textContent+=(l.textContent?'\n':'')+t;l.scrollTop=l.scrollHeight;}
function status(t,isErr){
  var b=document.getElementById('status');
  b.classList.add('on'); b.classList.toggle('err',!!isErr);
  document.getElementById('statusText').textContent=t;
  if(isErr) document.getElementById('log').classList.remove('hide');
}
function toggleLog(){
  var l=document.getElementById('log'), hidden=l.classList.toggle('hide');
  document.getElementById('toggleLog').textContent=hidden?'查看详情':'收起详情';
}
function up(){
  var f=document.getElementById('files').files;
  if(!f.length){alert('请先选择文件');return;}
  var kind=document.getElementById('kind').value,g=document.getElementById('grade').value;
  var i=0,added=0,dup=0,bad=0;
  (function next(){
    if(i>=f.length){
      log('全部完成。');
      status(bad
        ? ('完成，但有 '+bad+' 个文件未能解析；本次新增 '+added+' 条、重复忽略 '+dup+' 条')
        : ('解析完成：'+f.length+' 个文件，新增 '+added+' 条记录，重复忽略 '+dup+' 条'), bad>0);
      load(1);return;
    }
    var fd=new FormData();fd.append('file',f[i]);fd.append('kind',kind);fd.append('grade',g);
    status('正在解析 '+f[i].name+'（'+(i+1)+'/'+f.length+'）…');
    log('正在解析 '+f[i].name+' ...');
    fetch('api/upload',{method:'POST',body:fd}).then(r=>r.json()).then(d=>{
      if(d.error){bad++;log('  × '+d.error);}
      else{added+=d.n_new;dup+=d.n_dup;
           log('  √ 识别为['+d.kind_label+']，数据行 '+d.n_rows+
               '，新增证据 '+d.n_new+'，重复忽略 '+d.n_dup+(d.note?('，'+d.note):''));}
      i++;next();
    }).catch(e=>{bad++;log('  × '+e);i++;next();});
  })();
}
function load(p){
  if(p<1)p=1;page=p;
  var q=encodeURIComponent(document.getElementById('q').value||'');
  var sb=document.getElementById('sortBy').value;
  fetch('api/students?q='+q+'&page='+p+'&sort='+sb).then(r=>r.json()).then(d=>{
    var tb=document.getElementById('tb');
    if(!d.items.length){tb.innerHTML='<tr><td colspan="6" class="nodata">没有匹配的学生</td></tr>';}
    else{
      tb.innerHTML=d.items.map(function(s){
        var cells=s.grades.map(function(g){
          return g.has_data
            ? '<td><a class="btn sm" href="s/'+encodeURIComponent(s.sid)+'?grade='+encodeURIComponent(g.grade)+'">查看</a>'
              +'<div style="margin-top:3px;font-weight:800;color:#0f766e">'+g.cci.toFixed(2)+'</div></td>'
            : '<td><span class="nodata">未采集</span></td>';}).join('');
        return '<tr><td>'+s.sid+'</td><td>'+(s.name||'—')+'</td>'+cells+'</tr>';
      }).join('');
    }
    document.getElementById('cnt').textContent='共 '+d.total+' 名学生';
    document.getElementById('pg').textContent='第 '+d.page+' / '+d.pages+' 页';
  });
}
load(1);'''
        return page('画像总览', body, js)

    # ---------------- 上传接口 ----------------
    @app.post('/api/upload')
    async def api_upload(file: UploadFile = File(...), kind: str = Form('auto'),
                         grade: str = Form('大一')) -> JSONResponse:
        if grade not in GRADES:
            grade = GRADES[0]
        raw = await file.read()
        digest = sha256_bytes(raw)
        tmp = tempfile.NamedTemporaryFile(delete=False,
                                          suffix=os.path.splitext(file.filename)[1])
        tmp.write(raw)
        tmp.close()
        try:
            k = kind if kind != 'auto' else detect_kind(tmp.name)
            if k not in IMPORTERS:
                return JSONResponse({'error': f'无法识别文件类型：{file.filename}。'
                                              f'请在下拉框中手动指定材料类别。'})
            batch = sha256_text(digest, k, grade)[:16]
            res = IMPORTERS[k](STORE, tmp.name, grade, batch)
            STORE.log_upload(batch_id=batch, kind=k, filename=file.filename,
                             sha256=digest, grade=grade, **{
                                 key: res.get(key, 0) for key in ('n_rows', 'n_new', 'n_dup')},
                             note=res.get('note', ''))
            return JSONResponse({'kind': k, 'kind_label': KIND_LABEL.get(k, k), **res})
        except Exception as exc:                       # pragma: no cover
            return JSONResponse({'error': f'{type(exc).__name__}: {exc}'})
        finally:
            try:
                os.unlink(tmp.name)
            except Exception:
                pass

    @app.post('/api/upload/work')
    async def api_upload_work(sid: str = Form(...), grade: str = Form('大一'),
                              channel: str = Form('artifact'), note: str = Form(''),
                              file: UploadFile = File(...)) -> JSONResponse:
        if grade not in GRADES:
            grade = GRADES[0]
        if channel not in ('artifact', 'agent'):
            channel = 'artifact'
        if not STORE.conn().execute('SELECT 1 FROM students WHERE sid=?',
                                    (sid,)).fetchone():
            return JSONResponse({'error': f'学号 {sid} 不在学生库中，请先上传报名表'})
        raw = await file.read()
        tmp = tempfile.NamedTemporaryFile(
            delete=False, suffix=os.path.splitext(file.filename)[1])
        tmp.write(raw)
        tmp.close()
        try:
            batch = sha256_text(sha256_bytes(raw), channel, grade)[:16]
            res = import_artifact(STORE, tmp.name, file.filename, sid, grade,
                                  channel, batch, note)
            STORE.log_upload(batch_id=batch, kind=channel, filename=file.filename,
                             sha256=res['sha256'], grade=grade, n_rows=1,
                             n_new=res['n_new'], n_dup=res['n_dup'],
                             note=res.get('note', ''))
            return JSONResponse({'kind': channel,
                                 'kind_label': CHANNEL_LABEL[channel], **{
                                     k: v for k, v in res.items() if k != 'meta'},
                                 'meta': res.get('meta', {})})
        except Exception as exc:                       # pragma: no cover
            return JSONResponse({'error': f'{type(exc).__name__}: {exc}'})
        finally:
            try:
                os.unlink(tmp.name)
            except Exception:
                pass

    # ---------------- 学生列表 ----------------
    @app.get('/api/students')
    def api_students(q: str = Query(''), page_no: int = Query(1, alias='page'),
                     size: int = Query(20),
                     sort: str = Query('complete')) -> Dict[str, Any]:
        """学生列表。默认按材料完整度降序——先看结论最可信的那批学生。

        排序需要对候选集整体算完整度，因此先取全部匹配行再分页；
        学生量在万级以内，实测开销可接受。
        """
        cx = STORE.conn()
        size = max(1, min(200, size))
        where, args = '', []
        if q.strip():
            like = f'%{q.strip()}%'
            where = ('WHERE sid LIKE ? OR name LIKE ? OR class_name LIKE ? '
                     'OR college LIKE ? OR major LIKE ?')
            args = [like] * 5
        rows = cx.execute(
            f'SELECT * FROM students {where} ORDER BY class_name, sid', args).fetchall()
        total = len(rows)

        enriched = []
        for r in rows:
            pr = compute_portrait(STORE, r['sid'])
            comp = completeness_of(STORE, r['sid'], pr)
            best = max((pr[g]['cci'] for g in GRADES if pr[g]['has_data']), default=0.0)
            enriched.append((r, pr, comp, best))

        if sort == 'complete':
            enriched.sort(key=lambda x: (-x[2]['score'], -x[3]))
        elif sort == 'cci':
            enriched.sort(key=lambda x: (-x[3], -x[2]['score']))
        elif sort == 'cci_asc':
            enriched.sort(key=lambda x: (x[3], -x[2]['score']))
        elif sort == 'class':
            enriched.sort(key=lambda x: str(x[0]['sid']))

        pages = max(1, (total + size - 1) // size)
        page_no = max(1, min(page_no, pages))
        page_rows = enriched[(page_no - 1) * size: page_no * size]

        items = []
        for r, pr, comp, _best in page_rows:
            tag, cls = completeness_tag(comp['score'])
            items.append({
                'sid': r['sid'], 'name': r['name'], 'class_name': r['class_name'],
                'college': r['college'],
                'completeness': {**comp, 'tag': tag, 'cls': cls},
                'grades': [{'grade': g, 'has_data': pr[g]['has_data'],
                            'cci': pr[g]['cci']} for g in GRADES]})
        return {'total': total, 'page': page_no, 'pages': pages, 'items': items}

    @app.get('/api/portrait/{sid}')
    def api_portrait(sid: str) -> Dict[str, Any]:
        if not STORE.conn().execute('SELECT 1 FROM students WHERE sid=?',
                                    (sid,)).fetchone():
            raise HTTPException(404, '学生不存在')
        return {'student': student_row(STORE, sid),
                'grades': compute_portrait(STORE, sid)}

    @app.get('/api/report/{sid}')
    def api_report(sid: str, grade: str = Query('大一')) -> Response:
        if grade not in GRADES:
            raise HTTPException(400, '年级参数不合法')
        data = build_report_docx(STORE, sid, grade)
        name = f'{student_row(STORE, sid).get("name") or sid}-{grade}-文化德育画像报告.docx'
        quoted = name.encode('utf-8').hex()
        return Response(
            content=data,
            media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            headers={'Content-Disposition':
                     f"attachment; filename=report.docx; filename*=UTF-8''"
                     + ''.join(f'%{quoted[i:i + 2]}' for i in range(0, len(quoted), 2))})

    # ---------------- 学生详情页 ----------------
    @app.get('/s/{sid}', response_class=HTMLResponse)
    def detail(sid: str, grade: str = Query('大一')) -> str:
        if grade not in GRADES:
            grade = GRADES[0]
        if not STORE.conn().execute('SELECT 1 FROM students WHERE sid=?',
                                    (sid,)).fetchone():
            raise HTTPException(404, '学生不存在')
        stu = student_row(STORE, sid)
        all_g = compute_portrait(STORE, sid)
        pr = all_g[grade]
        lines = indicator_lines(STORE, sid, grade, limit_per_channel=3)
        matrix = dimension_matrix(STORE, sid, grade, pr)
        # 完整度只用于列表排序，不在界面上呈现；这里取值仅为统计证据渠道数
        comp = completeness_of(STORE, sid, all_g)
        advice = teacher_advice(pr, stu)

        tabs = ''.join(
            f'<a class="tab{" on" if g == grade else ""}" '
            f'href="?grade={g}">{g}{"" if all_g[g]["has_data"] else "（未采集）"}</a>'
            for g in GRADES)

        if lines:
            total_items = sum(ln['count'] for ln in lines)
            ind_html = (
                f'<div class="sum">本学年共采集到 <b>{total_items}</b> 项过程证据，'
                f'来自 <b>{comp["channels"]}</b> 类材料</div>'
                + ''.join(
                    f'<div class="ev"><b>{ln["label"]}</b> · {ln["count"]} 项'
                    + ''.join(f'<span class="it">· {it}</span>' for it in ln['items'])
                    + (f'<span class="it">· …… 另有 {ln["more"]} 项</span>'
                       if ln['more'] else '')
                    + '</div>' for ln in lines))
        else:
            ind_html = ('<div class="nodata">本学年尚未采集到过程性证据。'
                        '五维得分等于初始得分，请先上传该年级的导出表或学生作业。</div>')

        mat_html = ''.join(
            f'<div style="margin-bottom:11px">'
            f'<span class="pill p-{m["code"]}">{m["dim"]}</span> '
            f'<b style="color:#0f766e">{m["score"]:.2f}</b> '
            f'<span class="muted">初始得分 {m["baseline"]:.2f} → 证据充分度 '
            f'{m["rho"] * 100:.0f}%</span>'
            f'<div class="bar"><i style="width:{m["score"] * 100:.0f}%"></i></div>'
            f'<div class="muted" style="margin-top:4px">证据来源：'
            + ('；'.join(f'{s["name"]} {s["share"] * 100:.0f}%'
                         for s in m['sources'][:4]) or '暂无')
            + '</div></div>' for m in matrix)

        adv_html = ''.join(f'<div class="adv">{a}</div>' for a in advice)
        delta = pr['delta']
        growth = pr['growth'] * 100

        body = f'''
<div class="row" style="margin-bottom:14px">
  <a class="btn ghost sm" href="../">‹ 返回列表</a>
</div>
<div class="card">
  <div class="row" style="justify-content:space-between;align-items:flex-start">
    <div>
      <h2 style="margin-bottom:6px"><span class="n">◎</span>{stu.get('name') or '未命名'}
        · 文化德育数字画像</h2>
      <div class="muted">学号 {sid} ｜ 班级 {stu.get('class_name') or '—'}
        ｜ 学院 {stu.get('college') or '—'} ｜ 专业 {stu.get('major') or '—'}</div>

    </div>
    <div class="row">
      <a class="btn ghost" href="../api/report/{sid}?grade={grade}">生成学习报告</a>
      <a class="btn" href="../t/{sid}">查看成长轨迹</a>
    </div>
  </div>
  <div class="tabs" style="margin-top:14px">{tabs}</div>
</div>

<div class="grid2">
  <div class="card">
    <h2><span class="n">1</span>上传材料提取到的指标信息</h2>
    {ind_html}
  </div>
  <div class="card">
    <h2><span class="n">2</span>五维素养雷达（画像输出）</h2>
    <div class="rwrap">
      <div style="flex:1;min-width:260px">{radar_svg(pr['scores'], pr['baseline'])}</div>
      <div class="cci">
        <div style="font-size:15px;font-weight:700;color:#0f766e">综合指数</div>
        <div class="v">{pr['cci']:.2f}</div>
        <div style="font-size:14.5px;color:#047857;font-weight:800">
          {'↑' if delta >= 0 else '↓'} 较上学年 {delta:+.2f}</div>
        <div style="font-size:14px;font-weight:600;color:#5b6b7c">环比增长 {growth:+.1f}%</div>
        <div style="font-size:15px;font-weight:800;color:#065f46;margin-top:7px">
          {level_of(pr['cci'])}</div>
      </div>
    </div>
    <div class="radarnote">实线为当前得分，虚线为初始得分。</div>
  </div>
  <div class="card">
    <h2><span class="n">3</span>指标与五维的对应关系（量化依据）</h2>
    {mat_html}
    <div class="muted" style="margin-top:8px">评分可溯源：每一项得分均可回溯到上表列出的
      具体学习事件与过程数据。</div>
  </div>
  <div class="card">
    <h2><span class="n">4</span>给授课教师的改进建议</h2>
    {adv_html}
  </div>
</div>

<div class="card">
  <h2><span class="n">5</span>补充上传该生的作业成果 / 自建智能体</h2>
  <div class="up">
    <div class="row">
      <select id="ch"><option value="artifact">作业成果（视频 / Word / Excel / PPT / PDF）</option>
        <option value="agent">学生自建智能体（提示词、知识库导出文件）</option></select>
      <select id="g2">{''.join(f'<option{" selected" if g == grade else ""}>{g}</option>'
                               for g in GRADES)}</select>
      <input type="file" id="wf" multiple>
      <button class="btn" onclick="upw()">上传并计分</button>
    </div>
    <input id="note" placeholder="可选：补充说明（作品主题、创作思路等，将参与文化主题识别）"
           style="width:100%;margin-top:9px">
  </div>
  <div class="statusbar" id="wstatus">
    <span class="dot"></span><span id="wstatusText"></span>
    <span class="lnk" id="wToggle" onclick="wToggleLog()">查看详情</span>
  </div>
  <div class="log hide" id="wlog"></div>
</div>'''

        js = f'''
var SID={json.dumps(sid)};
function wlog(t){{var l=document.getElementById('wlog');
  l.textContent+=(l.textContent?'\\n':'')+t;l.scrollTop=l.scrollHeight;}}
function wstatus(t,isErr){{
  var b=document.getElementById('wstatus');
  b.classList.add('on'); b.classList.toggle('err',!!isErr);
  document.getElementById('wstatusText').textContent=t;
  if(isErr) document.getElementById('wlog').classList.remove('hide');
}}
function wToggleLog(){{
  var l=document.getElementById('wlog'), hidden=l.classList.toggle('hide');
  document.getElementById('wToggle').textContent=hidden?'查看详情':'收起详情';
}}
function upw(){{
  var f=document.getElementById('wf').files;
  if(!f.length){{alert('请先选择文件');return;}}
  var ch=document.getElementById('ch').value,g=document.getElementById('g2').value;
  var note=document.getElementById('note').value||'';
  var i=0,added=0,dup=0,bad=0;
  (function next(){{
    if(i>=f.length){{
      wlog('全部完成。');
      wstatus(bad
        ? ('完成，但有 '+bad+' 个文件未能解析；已计入 '+added+' 份')
        : ('已计入 '+added+' 份材料'+(dup?('，'+dup+' 份为重复上传未重复计分'):'')+'，正在刷新画像…'), bad>0);
      if(!bad) setTimeout(function(){{
        location.href='?grade='+encodeURIComponent(g);}},900);
      return;
    }}
    var fd=new FormData();
    fd.append('file',f[i]);fd.append('sid',SID);fd.append('grade',g);
    fd.append('channel',ch);fd.append('note',note);
    wstatus('正在解析 '+f[i].name+'（'+(i+1)+'/'+f.length+'）…');
    wlog('正在解析 '+f[i].name+' ...');
    fetch('../api/upload/work',{{method:'POST',body:fd}}).then(r=>r.json()).then(d=>{{
      if(d.error){{bad++;wlog('  × '+d.error);}}
      else{{d.n_new?added++:dup++;
           wlog('  '+(d.n_new?'√':'·')+' '+d.kind_label+'：'+(d.note||'')+
               (d.keywords&&d.keywords.length?('｜主题词 '+d.keywords.slice(0,6).join('、')):''));}}
      i++;next();
    }}).catch(e=>{{bad++;wlog('  × '+e);i++;next();}});
  }})();
}}'''
        return page(f'{stu.get("name") or sid} · {grade}', body, js)

    # ---------------- 成长轨迹页 ----------------
    @app.get('/t/{sid}', response_class=HTMLResponse)
    def trajectory(sid: str) -> str:
        if not STORE.conn().execute('SELECT 1 FROM students WHERE sid=?',
                                    (sid,)).fetchone():
            raise HTTPException(404, '学生不存在')
        stu = student_row(STORE, sid)
        all_g = compute_portrait(STORE, sid)

        head = ''.join(f'<th>{g}</th>' for g in GRADES)
        rows = ''
        for d in DIMS:
            cells = ''.join(
                (f'<td><b>{all_g[g]["scores"][d]:.2f}</b></td>' if all_g[g]['has_data']
                 else '<td class="nodata">—</td>') for g in GRADES)
            rows += (f'<tr><td style="text-align:left">'
                     f'<span class="pill p-{DIM_CODE[d]}">{d}</span></td>{cells}</tr>')
        cci_cells = ''.join(
            (f'<td><b style="color:#065f46">{all_g[g]["cci"]:.2f}</b>'
             f'<div class="muted">{level_of(all_g[g]["cci"])}</div></td>'
             if all_g[g]['has_data'] else '<td class="nodata">—</td>') for g in GRADES)
        rows += f'<tr><td style="text-align:left"><b>综合指数 CCI</b></td>{cci_cells}</tr>'

        body = f'''
<div class="row" style="margin-bottom:14px">
  <a class="btn ghost sm" href="../s/{sid}">‹ 返回画像</a>
</div>
<div class="card">
  <h2><span class="n">◎</span>{stu.get('name') or sid} · 四学年德育成长轨迹</h2>
  <div class="muted" style="margin-bottom:10px">学号 {sid} ｜ 班级
    {stu.get('class_name') or '—'}。虚线为综合指数，各学年初始得分自动承接上一学年结果，
    因此曲线体现的是德育养成的累积过程。</div>
  {trajectory_svg(all_g)}
</div>
<div class="card">
  <h2><span class="n">◎</span>逐年得分表</h2>
  <table><thead><tr><th style="text-align:left">维度</th>{head}</tr></thead>
  <tbody>{rows}</tbody></table>
</div>'''
        return page(f'{stu.get("name") or sid} · 成长轨迹', body)

    @app.get('/api/health')
    def health() -> Dict[str, Any]:
        cx = STORE.conn()
        return {
            'app': APP_NAME, 'version': APP_VERSION, 'db': STORE.db_path,
            'students': cx.execute('SELECT COUNT(*) n FROM students').fetchone()['n'],
            'evidences': cx.execute('SELECT COUNT(*) n FROM evidences').fetchone()['n'],
            'baselines': cx.execute(
                'SELECT COUNT(DISTINCT sid) n FROM baselines').fetchone()['n'],
            'uploads': cx.execute('SELECT COUNT(*) n FROM uploads').fetchone()['n'],
        }
else:                                                   # pragma: no cover
    app = None


# ==========================================================================
# 十一、命令行入口
# ==========================================================================
def _cli() -> int:
    import argparse
    ap = argparse.ArgumentParser(description=f'{APP_NAME} v{APP_VERSION}')
    sub = ap.add_subparsers(dest='cmd')

    p_imp = sub.add_parser('import', help='导入一个导出表（自动识别类型）')
    p_imp.add_argument('path')
    p_imp.add_argument('--grade', default='大一', choices=list(GRADES))
    p_imp.add_argument('--kind', default='auto')

    p_sc = sub.add_parser('score', help='打印某学生的画像')
    p_sc.add_argument('sid')

    p_srv = sub.add_parser('serve', help='启动 Web 服务')
    p_srv.add_argument('--host', default='127.0.0.1')
    p_srv.add_argument('--port', type=int, default=8000)

    sub.add_parser('health', help='查看数据库统计')
    sub.add_parser('reset-evidence',
                   help='清空过程证据（保留学生库与问卷初始得分），用于调整计量参数后重新导入')

    args = ap.parse_args()
    if args.cmd == 'import':
        kind = args.kind if args.kind != 'auto' else detect_kind(args.path)
        if kind not in IMPORTERS:
            print(f'无法识别文件类型：{args.path}')
            return 2
        batch = sha256_text(args.path, kind, args.grade)[:16]
        res = IMPORTERS[kind](STORE, args.path, args.grade, batch)
        print(f'[{KIND_LABEL.get(kind, kind)}] {res}')
    elif args.cmd == 'score':
        pr = compute_portrait(STORE, args.sid)
        print(json.dumps({'student': student_row(STORE, args.sid), 'grades': pr},
                         ensure_ascii=False, indent=2))
    elif args.cmd == 'reset-evidence':
        cx = STORE.conn()
        n = cx.execute('SELECT COUNT(*) n FROM evidences').fetchone()['n']
        cx.execute('DELETE FROM evidences')
        cx.execute('DELETE FROM uploads')
        cx.commit()
        print(f'已清空 {n} 条证据与上传记录，学生库与问卷初始得分保留。请重新导入各表。')
    elif args.cmd == 'health':
        print(json.dumps(health(), ensure_ascii=False, indent=2))
    elif args.cmd == 'serve':
        import uvicorn
        uvicorn.run(app, host=args.host, port=args.port)
    else:
        ap.print_help()
    return 0


if __name__ == '__main__':
    raise SystemExit(_cli())
