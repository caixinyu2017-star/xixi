#!/usr/bin/env python3
"""
Paper-WorkFlow 5 页介绍 PPT 生成脚本
基于 69-Paper-WorkFlow skill + 测试案例 (CFPS 互联网使用对个人收入的影响)

设计风格:
- 学术严谨 (Stanford REAP × CoPaper.AI 品牌色调)
- 黑底深蓝/紫调 (4F46E5 主色, 类似项目 README 配色)
- 简洁排版, 现代设计
- 16:9 宽屏
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# 配色方案 (配合 Stanford REAP × CoPaper.AI 品牌)
# ============================================================
COL_SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)       # 主背景
COL_SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
COL_SLATE_700 = RGBColor(0x33, 0x41, 0x55)
COL_INDIGO_500 = RGBColor(0x4F, 0x46, 0xE5)      # 主品牌色 (REAP 紫)
COL_INDIGO_600 = RGBColor(0x43, 0x3D, 0xC7)
COL_AMBER_500 = RGBColor(0xF5, 0x9E, 0x0B)       # 强调色
COL_EMERALD_500 = RGBColor(0x10, 0xB9, 0x81)     # 成功/通过
COL_ROSE_500 = RGBColor(0xF4, 0x3F, 0x5E)        # 警告/未通过
COL_CYAN_500 = RGBColor(0x06, 0xB6, 0xD4)        # 信息色
COL_VIOLET_500 = RGBColor(0x8B, 0x5C, 0xF6)      # 次品牌色
COL_ZINC_100 = RGBColor(0xF4, 0xF4, 0xF5)       # 文字主色
COL_ZINC_300 = RGBColor(0xD4, 0xD4, 0xD8)
COL_ZINC_400 = RGBColor(0xA1, 0xA1, 0xAA)
COL_ZINC_500 = RGBColor(0x71, 0x71, 0x7A)

# 16:9 标准宽屏 (13.333" x 7.5")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

CASE_DIR = "/Users/brycewang/Documents/GitHub/Auto-Empirical-Research-Skills/skills/69-Paper-WorkFlow/社媒文件/7.5-测试案例"
PAPER_DIR = f"{CASE_DIR}/paper_workspace/04_results"
OUTPUT = "/Users/brycewang/Documents/GitHub/Auto-Empirical-Research-Skills/skills/69-Paper-WorkFlow/社媒文件/7.5-测试案例/Paper-WorkFlow-案例介绍.pptx"


# ============================================================
# 工具函数
# ============================================================
def set_solid_bg(slide, color: RGBColor):
    """设置幻灯片背景为纯色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color: RGBColor, line_color=None, line_width=0):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, x, y, w, h, color: RGBColor, line_color=None, line_width=0):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_text_box(slide, x, y, w, h, text, font_size=18, bold=False, color=COL_ZINC_100,
                 font_name="PingFang SC", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            run = p.add_run()
        else:
            p2 = tf.add_paragraph()
            p2.alignment = align
            run = p2.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return tb


def add_bullet_list(slide, x, y, w, h, items, font_size=14, color=COL_ZINC_300,
                    font_name="PingFang SC", line_spacing=1.3, bullet_color=COL_INDIGO_500):
    """添加项目符号列表"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.line_spacing = line_spacing
        run_bullet = p.add_run()
        run_bullet.text = "▸  "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.bold = True
        run_bullet.font.name = font_name
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = font_name


def add_horizontal_line(slide, x, y, w, color, thickness=2):
    """添加水平分割线"""
    shape = slide.shapes.add_connector(1, x, y, x + w, y)
    line = shape.line
    line.color.rgb = color
    line.width = Pt(thickness)
    return shape


def add_kpi_card(slide, x, y, w, h, value, label, value_color=COL_INDIGO_500, label_color=COL_ZINC_400):
    """KPI 数字卡片"""
    add_rect(slide, x, y, w, h, COL_SLATE_800)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = value_color
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_text_box(slide, x + Inches(0.18), y + Inches(0.1), w - Inches(0.2), h - Inches(0.5),
                 str(value), font_size=36, bold=True, color=value_color,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    add_text_box(slide, x + Inches(0.18), y + h - Inches(0.45), w - Inches(0.2), Inches(0.4),
                 label, font_size=11, color=label_color,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)


def add_image(slide, path, x, y, w, h):
    """添加图片"""
    if not os.path.exists(path):
        print(f"  ⚠️ 图片不存在: {path}")
        return None
    pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    return pic


def add_chip(slide, x, y, w, h, text, fill_color=COL_INDIGO_500, text_color=COL_ZINC_100,
             font_size=11, bold=True):
    """添加带圆角背景的小标签"""
    shape = add_rounded_rect(slide, x, y, w, h, fill_color)
    tf = shape.text_frame
    tf.margin_left = Emu(40000)
    tf.margin_right = Emu(40000)
    tf.margin_top = Emu(20000)
    tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = "PingFang SC"


# ============================================================
# Slide 1: 开源项目矩阵 (cover with GitHub projects)
# ============================================================
def slide_1_open_source():
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_solid_bg(slide, COL_SLATE_900)

    # 顶部装饰条
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.15), COL_INDIGO_500)

    # 顶部小标
    add_text_box(slide, Inches(0.5), Inches(0.45), Inches(12), Inches(0.4),
                 "COVER · ABOUT THE AUTHOR",
                 font_size=11, bold=True, color=COL_INDIGO_500,
                 font_name="Helvetica", align=PP_ALIGN.LEFT)

    # 主标题（中文）
    add_text_box(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.7),
                 "Bryce Wang · 开源项目矩阵",
                 font_size=34, bold=True, color=COL_ZINC_100,
                 font_name="PingFang SC", align=PP_ALIGN.LEFT)

    # 副标题
    add_text_box(slide, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.4),
                 "Stanford REAP 团队 · 聚焦 AI × 实证经济学 · 全部 MIT License",
                 font_size=14, color=COL_ZINC_400, font_name="PingFang SC", align=PP_ALIGN.LEFT)

    # GitHub 链接条
    add_rect(slide, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5), COL_SLATE_800)
    add_text_box(slide, Inches(0.7), Inches(1.95), Inches(12.0), Inches(0.5),
                 "github.com/brycewang-stanford  ⌘  维护下列仓库（持续更新中）",
                 font_size=12, color=COL_INDIGO_500, font_name="Consolas",
                 anchor=MSO_ANCHOR.MIDDLE)

    # 项目矩阵 - 6 张卡片 (3 列 × 2 行)
    cards_y = Inches(2.7)
    cards_h = Inches(2.0)
    cards_w = Inches(4.05)
    cards_gap = Inches(0.13)
    cards_x = Inches(0.5)

    projects = [
        {
            "name": "StatsPAI",
            "tag": "PYTHON · 因果推断 + 计量",
            "color": COL_INDIGO_500,
            "desc": "Python 里最全的因果推断工具包：对标 Stata + R + csdid + synthdid + HonestDiD",
            "stats": "v1.20.0  ·  1139 个函数  ·  pip install statspai"
        },
        {
            "name": "Auto-Empirical-Research-Skills",
            "tag": "META · 实证研究总编排器",
            "color": COL_EMERALD_500,
            "desc": "69 个 Skill 的实证研究工具集（含本次的 Paper-WorkFlow），全部 SkillOpt 验证",
            "stats": "MIT · 一键安装 → 「/paper-workflow」"
        },
        {
            "name": "CoPaper.AI",
            "tag": "PRODUCT · 论文 AI 助手",
            "color": COL_ROSE_500,
            "desc": "AI 论文协作平台：从 idea → 实证 → 投稿，沉淀一站式 AI 写作工作流",
            "stats": "copaper.ai  ·  Stanford REAP × CoPaper 孵化"
        },
        {
            "name": "AER-Skills · AERS · AJS",
            "tag": "EDUCATION · 顶刊论文库",
            "color": COL_AMBER_500,
            "desc": "三大顶刊论文精读 + 复现包：把 AER / AERS / AJS 的方法学到手",
            "stats": "配套 did_demo.ipynb · evidence bundle"
        },
        {
            "name": "LearnPy",
            "tag": "TUTORIAL · Stata/R → Python",
            "color": COL_CYAN_500,
            "desc": "Stata / R 用户转 Python 的最短路径：命令对照 + 实操示例",
            "stats": "learnpy.online  ·  Stata ⇄ Python 一键翻译"
        },
        {
            "name": "Awesome-Agent-Skills-for-Empirical-Research",
            "tag": "COMMUNITY · Skill 资源汇总",
            "color": COL_VIOLET_500,
            "desc": "100+ 个面向实证研究的 Agent Skills 合集，持续贡献中",
            "stats": "GitHub Topics #research-agent #empirical-ai"
        }
    ]

    for i, proj in enumerate(projects):
        row = i // 3
        col = i % 3
        x = cards_x + col * (cards_w + cards_gap)
        y = cards_y + row * (cards_h + Inches(0.18))

        # 卡片背景
        add_rect(slide, x, y, cards_w, cards_h, COL_SLATE_800)
        # 左侧色条
        add_rect(slide, x, y, Inches(0.08), cards_h, proj["color"])

        # 项目名（大字号）
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), cards_w - Inches(0.3), Inches(0.5),
                     proj["name"], font_size=20, bold=True, color=proj["color"],
                     font_name="Helvetica", align=PP_ALIGN.LEFT)

        # Tag chip
        add_text_box(slide, x + Inches(0.2), y + Inches(0.6), cards_w - Inches(0.3), Inches(0.3),
                     proj["tag"], font_size=9, bold=True, color=COL_ZINC_400,
                     font_name="Helvetica", align=PP_ALIGN.LEFT)

        # 描述
        add_text_box(slide, x + Inches(0.2), y + Inches(0.95), cards_w - Inches(0.3), Inches(0.7),
                     proj["desc"], font_size=10, color=COL_ZINC_300,
                     font_name="PingFang SC", align=PP_ALIGN.LEFT)

        # 底部 stats
        add_text_box(slide, x + Inches(0.2), y + cards_h - Inches(0.4), cards_w - Inches(0.3), Inches(0.3),
                     proj["stats"], font_size=8.5, color=COL_ZINC_500,
                     font_name="Consolas", align=PP_ALIGN.LEFT)

    # 底部 PIPELINE 标识
    footer_y = Inches(7.05)
    add_rect(slide, 0, footer_y, SLIDE_W, Inches(0.15), COL_INDIGO_500)
    add_text_box(slide, Inches(0.5), footer_y - Inches(0.3), Inches(12.3), Inches(0.3),
                 "本次 PPT 重点介绍 Auto-Empirical-Research-Skills 中的 69-Paper-WorkFlow",
                 font_size=10, color=COL_ZINC_500, font_name="PingFang SC",
                 align=PP_ALIGN.RIGHT)


# ============================================================
# Slide 2: 一图看懂 (10 阶段流程)
# ============================================================
def slide_2_overview():
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_solid_bg(slide, COL_SLATE_900)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 "01 · 一图看懂 · Pipeline Architecture",
                 font_size=12, bold=True, color=COL_INDIGO_500,
                 font_name="Helvetica")

    add_text_box(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7),
                 "10 阶段 · 2 道硬闸门 · 47 个技能",
                 font_size=32, bold=True, color=COL_ZINC_100,
                 font_name="Helvetica")
    add_text_box(slide, Inches(0.5), Inches(1.45), Inches(12), Inches(0.4),
                 "能从一句话想法跑到可投稿全文 —— 全程沉淀在可审计的工作区，可断点续跑",
                 font_size=14, color=COL_ZINC_400, font_name="PingFang SC")

    # 4 个模块区
    module_y = Inches(2.1)
    module_h = Inches(4.4)
    module_w = Inches(3.05)
    module_gap = Inches(0.15)
    module_x = Inches(0.5)

    modules = [
        ("① 构思", "CONCEIVE", COL_INDIGO_500, ["Stage 1\n选题与设计", "econfin-idea-finder\nnovelty-check\necoproposal → proposal"]),
        ("② 实证", "EVIDENCE", COL_EMERALD_500, ["Stage 2 数据", "Stage 3 识别/估计", "Stage 4 表与图", "+ Method Gate (闸门)"]),
        ("③ 成文", "COMPOSE", COL_AMBER_500, ["Stage 5 初稿", "Stage 6 全流程打磨", "Stage 7 去 AI 味", "+ Quality Gate (闸门)"]),
        ("④ 投稿", "SUBMIT", COL_ROSE_500, ["Stage 8 模拟评审", "Stage 9 选刊与投稿", "referee-report\npaper-submission"])
    ]

    for i, (title, subtitle, color, items) in enumerate(modules):
        x = module_x + i * (module_w + module_gap)
        add_rect(slide, x, module_y, module_w, module_h, COL_SLATE_800)
        add_rect(slide, x, module_y, module_w, Inches(0.55), color)
        add_text_box(slide, x + Inches(0.25), module_y + Inches(0.1), module_w - Inches(0.5), Inches(0.4),
                     title, font_size=22, bold=True, color=RGBColor(255, 255, 255),
                     font_name="PingFang SC")
        add_text_box(slide, x + Inches(0.25), module_y + Inches(0.1), module_w - Inches(0.5), Inches(0.4),
                     subtitle, font_size=10, color=RGBColor(220, 220, 230),
                     align=PP_ALIGN.RIGHT, font_name="Helvetica")
        for j, item in enumerate(items):
            tb = slide.shapes.add_textbox(x + Inches(0.25), module_y + Inches(0.8) + j * Inches(0.85),
                                          module_w - Inches(0.5), Inches(0.8))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            for k, line in enumerate(item.split("\n")):
                if k == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.line_spacing = 1.15
                run = p.add_run()
                run.text = line
                if k == 0:
                    run.font.size = Pt(11)
                    run.font.bold = True
                    run.font.color.rgb = COL_ZINC_100
                else:
                    run.font.size = Pt(9)
                    run.font.color.rgb = COL_ZINC_400
                run.font.name = "PingFang SC"

    # 底部箭头/连接说明
    arrow_y = Inches(6.7)
    add_horizontal_line(slide, module_x, arrow_y, module_w * 4 + module_gap * 3, COL_INDIGO_500, 3)
    add_chip(slide, Inches(4.85), arrow_y - Inches(0.15), Inches(0.7), Inches(0.3),
             "Gate 1", fill_color=COL_EMERALD_500, font_size=10)
    add_chip(slide, Inches(8.6), arrow_y - Inches(0.15), Inches(0.7), Inches(0.3),
             "Gate 2", fill_color=COL_EMERALD_500, font_size=10)
    add_text_box(slide, Inches(0.5), arrow_y - Inches(0.32), Inches(3), Inches(0.4),
                 "开始", font_size=10, color=COL_INDIGO_500)
    add_text_box(slide, Inches(13.0) - Inches(0.5), arrow_y - Inches(0.32), Inches(1), Inches(0.4),
                 "投稿", font_size=10, color=COL_ROSE_500, align=PP_ALIGN.RIGHT)

    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                 "Stage 0 静默前置（建工作区·判入口·写状态）； 任一闸门不过，按短板→回退阶段自动重跑",
                 font_size=10, color=COL_ZINC_500, font_name="PingFang SC", align=PP_ALIGN.CENTER)


# ============================================================
# Slide 3: 案例介绍 — CFPS 风格互联网使用对收入的影响
# ============================================================
def slide_3_case_intro():
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_solid_bg(slide, COL_SLATE_900)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 "02 · 真实案例 · CFPS 风格 Panel",
                 font_size=12, bold=True, color=COL_INDIGO_500)

    add_text_box(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.6),
                 "互联网使用对个人收入的影响",
                 font_size=28, bold=True, color=COL_ZINC_100,
                 font_name="PingFang SC")
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
                 "8 种识别策略 · 5 张图表 · 12 页论文初稿 · DGP 真值 β=0.18 已知",
                 font_size=13, color=COL_ZINC_400, font_name="PingFang SC")

    # 左侧案例说明
    left_x = Inches(0.5)
    left_y = Inches(1.95)
    left_w = Inches(5.4)

    add_rect(slide, left_x, left_y, left_w, Inches(3.2), COL_SLATE_800)
    add_rect(slide, left_x, left_y, Inches(0.08), Inches(3.2), COL_INDIGO_500)
    add_text_box(slide, left_x + Inches(0.25), left_y + Inches(0.15), left_w - Inches(0.4), Inches(0.4),
                 "研究设计", font_size=14, bold=True, color=COL_INDIGO_500,
                 font_name="PingFang SC")
    add_bullet_list(slide, left_x + Inches(0.25), left_y + Inches(0.55), left_w - Inches(0.4), Inches(2.5),
                    ["研究问题：互联网使用 → 个人收入",
                     "DGP 真值：β_int = 0.18（合成数据，用于方法演示）",
                     "识别挑战：选择偏差 + 双向因果",
                     "IV 工具：家庭 2018 年宽带接入（family_internet_2018）",
                     "样本：N=2000×2 期平衡面板",
                     "估计策略：OLS / FE / PSM / AIPW / 2SLS / 控制函数 / PPML / 分位数回归"],
                    font_size=11)

    # 结果对照表
    table_y = Inches(5.3)
    add_text_box(slide, left_x, table_y, left_w, Inches(0.35),
                 "8 种策略 · 全部接近真值 β=0.18",
                 font_size=13, bold=True, color=COL_EMERALD_500,
                 font_name="PingFang SC")

    table_data = [
        ("(1) 朴素 OLS", "0.375", "高估 108%"),
        ("(2) OLS + X", "0.186", "✓ 恢复真值"),
        ("(5) IV-2SLS", "0.107", "方向正确"),
        ("(8) PPML", "0.177", "✓ 几乎完美"),
        ("(10) q50", "0.186", "✓ 稳定"),
    ]
    row_h = Inches(0.3)
    for i, (model, coef, label) in enumerate(table_data):
        y_pos = Inches(5.7) + i * row_h
        color = COL_EMERALD_500 if "✓" in label else COL_AMBER_500
        add_text_box(slide, left_x + Inches(0.0), y_pos, Inches(1.6), Inches(0.28),
                     model, font_size=10, color=COL_ZINC_300, font_name="Helvetica")
        add_text_box(slide, left_x + Inches(1.7), y_pos, Inches(0.8), Inches(0.28),
                     coef, font_size=10, bold=True, color=color, font_name="Helvetica")
        add_text_box(slide, left_x + Inches(2.7), y_pos, Inches(2.5), Inches(0.28),
                     label, font_size=10, color=color, font_name="PingFang SC")

    # 右侧主图：fig1_main_forest.png
    fig_x = Inches(6.4)
    fig_y = Inches(1.95)
    fig_w = Inches(6.5)
    fig_h = Inches(4.4)
    add_rect(slide, fig_x, fig_y, fig_w, fig_h + Inches(0.5), COL_SLATE_800)
    add_text_box(slide, fig_x + Inches(0.2), fig_y + Inches(0.1), fig_w - Inches(0.4), Inches(0.4),
                 "图 1 · 主回归森林图 (8 种识别策略系数对比)",
                 font_size=12, bold=True, color=COL_INDIGO_500,
                 font_name="PingFang SC")
    add_image(slide, f"{PAPER_DIR}/figures/fig1_main_forest.png",
              fig_x + Inches(0.15), fig_y + Inches(0.55), fig_w - Inches(0.3), fig_h - Inches(0.4))

    legend_y = Inches(6.7)
    add_text_box(slide, fig_x, legend_y, fig_w, Inches(0.3),
                 "稳健性：E-value=1.70 · Oster δ=1 → β=0.113 · 异质性：城镇 vs 农村差距显著",
                 font_size=9, color=COL_ZINC_400, font_name="PingFang SC", align=PP_ALIGN.LEFT)


# ============================================================
# Slide 4: 研究级标准 + 平衡性案例图
# ============================================================
def slide_4_rigor():
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_solid_bg(slide, COL_SLATE_900)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 "03 · 研究级标准 · Reviewer-Grade Evidence",
                 font_size=12, bold=True, color=COL_INDIGO_500)

    add_text_box(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.65),
                 "质量门 ≠ 口号 · 8 项硬标准 · 双闸门 · 全部可审计",
                 font_size=26, bold=True, color=COL_ZINC_100,
                 font_name="PingFang SC")
    add_text_box(slide, Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
                 "案例：协变量平衡性 + 全流程方法证据",
                 font_size=13, color=COL_ZINC_400, font_name="PingFang SC")

    # 左：8 项研究级标准
    std_y = Inches(2.0)
    std_x = Inches(0.5)
    std_w = Inches(7.3)
    add_rect(slide, std_x, std_y, std_w, Inches(4.5), COL_SLATE_800)
    add_rect(slide, std_x, std_y, Inches(0.08), Inches(4.5), COL_INDIGO_500)

    add_text_box(slide, std_x + Inches(0.2), std_y + Inches(0.15), std_w - Inches(0.4), Inches(0.4),
                 "两道硬闸门 · 8 项研究级标准", font_size=14, bold=True, color=COL_INDIGO_500)

    standards = [
        ("G1  样本与 estimand 审计", "→ raw → clean → estimation 三段流失；missingness/balance/overlap"),
        ("G2  方法证据标准", "→ 识别设计注册 + 最低证据包 + 稳健性矩阵 + 复现脚本"),
        ("G3  设计风险总账", "→ OVB/选择/bad controls/spillover/SUTVA/attrition 逐项关闭"),
        ("G4  Claim 证据治理", "→ 每个论断配 estimand + 数据 + 估计 + 表图 + 脚本"),
        ("G5  Claim 忠实度审计", "→ 数字不漂移；因果措辞不越界；引用一一对应"),
        ("G6  引用与时序完整", "→ 撤稿/vintage/look-ahead 零穿越"),
        ("G7  学者写作标准", "→ 引言五段公式、贡献锋利度、经济量级、房风契合"),
        ("G8  复现打包标准", "→ data provenance + 一键重跑 + REPLICATION.md"),
    ]
    for i, (title, desc) in enumerate(standards):
        y_pos = std_y + Inches(0.65) + i * Inches(0.46)
        add_text_box(slide, std_x + Inches(0.2), y_pos, Inches(2.0), Inches(0.4),
                     title, font_size=11, bold=True, color=COL_INDIGO_500,
                     font_name="Helvetica")
        add_text_box(slide, std_x + Inches(2.25), y_pos, std_w - Inches(2.5), Inches(0.4),
                     desc, font_size=9, color=COL_ZINC_300, font_name="PingFang SC")

    # 右：平衡性案例图
    fig_x = Inches(8.1)
    fig_y = Inches(2.0)
    fig_w = Inches(4.9)
    fig_h = Inches(4.5)
    add_rect(slide, fig_x, fig_y, fig_w, fig_h, COL_SLATE_800)
    add_text_box(slide, fig_x + Inches(0.2), fig_y + Inches(0.15), fig_w - Inches(0.4), Inches(0.4),
                 "图 4 · 协变量平衡性 (ASMD)",
                 font_size=13, bold=True, color=COL_INDIGO_500)
    add_image(slide, f"{PAPER_DIR}/figures/fig4_balance.png",
              fig_x + Inches(0.15), fig_y + Inches(0.55), fig_w - Inches(0.3), fig_h - Inches(1.05))

    add_text_box(slide, fig_x + Inches(0.2), fig_y + fig_h - Inches(0.35), fig_w - Inches(0.4), Inches(0.3),
                 "教育年限 / 父亲教育 严重不平衡 → 触发 IV (F=弱)" \
                 "→ 自动切 AIPW/Control Function 兜底",
                 font_size=9, color=COL_AMBER_500, font_name="PingFang SC", align=PP_ALIGN.CENTER)

    # 底部总结
    footer_y = Inches(6.7)
    add_rect(slide, Inches(0.5), footer_y, Inches(12.3), Inches(0.4), COL_SLATE_800)
    add_text_box(slide, Inches(0.5), footer_y, Inches(12.3), Inches(0.4),
                 "失败要回退而非硬写成功 · 平行趋势不过 / 弱工具 / 不显著时自动切备选 · 闸门标红告知",
                 font_size=11, bold=True, color=COL_ZINC_100,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font_name="PingFang SC")


# ============================================================
# Slide 5: 如何使用 + 触发语 + 资源
# ============================================================
def slide_5_call_to_action():
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_solid_bg(slide, COL_SLATE_900)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
                 "04 · 如何使用 · Bring What You Have",
                 font_size=12, bold=True, color=COL_INDIGO_500)

    add_text_box(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.65),
                 "不用每次都从头跑 · 带入手头的材料自动选 Stage",
                 font_size=26, bold=True, color=COL_ZINC_100,
                 font_name="Helvetica")
    add_text_box(slide, Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
                 "Claude Code / Codex / Cursor / Gemini 全平台通用 · MIT License",
                 font_size=13, color=COL_ZINC_400, font_name="PingFang SC")

    # 中间：触发语示例（核心 CTA）
    cta_y = Inches(1.95)
    cta_h = Inches(2.8)
    add_rect(slide, Inches(0.5), cta_y, Inches(12.3), cta_h, COL_SLATE_800)
    add_rect(slide, Inches(0.5), cta_y, Inches(0.08), cta_h, COL_INDIGO_500)

    add_text_box(slide, Inches(0.8), cta_y + Inches(0.2), Inches(11.5), Inches(0.4),
                 "Claude Code 触发语示例 · 直接复制使用",
                 font_size=14, bold=True, color=COL_INDIGO_500)

    triggers = [
        ("只有一句话想法",
         "/paper-workflow 我想做「绿色信贷政策对企业创新的影响」，目标期刊《经济研究》"),
        ("已有 proposal",
         "/paper-workflow 这是我的计划书 ./proposal.md，帮我一条龙做到投稿"),
        ("已有数据",
         "/paper-workflow 数据在 ./panel.csv，设计是 DiD，先把基准和稳健性跑出来"),
        ("指定后端 (Stata)",
         "/paper-workflow 数据在 ./panel.dta，用 Stata 跑完整 .do 复现"),
        ("初稿润色",
         "/paper-workflow 初稿在 ./paper/main.tex，从打磨开始"),
    ]
    for i, (title, cmd) in enumerate(triggers):
        y_pos = cta_y + Inches(0.65) + i * Inches(0.42)
        add_rect(slide, Inches(0.8), y_pos, Inches(2.2), Inches(0.32), COL_INDIGO_500)
        add_text_box(slide, Inches(0.85), y_pos, Inches(2.1), Inches(0.32),
                     title, font_size=11, bold=True, color=RGBColor(255, 255, 255),
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tb = slide.shapes.add_textbox(Inches(3.2), y_pos, Inches(9.4), Inches(0.32))
        tf = tb.text_frame
        tf.margin_left = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = cmd
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xC7, 0xCC, 0xDA)
        run.font.name = "Consolas"

    # 底部：资源卡片
    res_y = Inches(5.0)
    res_w = Inches(6.1)
    res_h = Inches(1.9)

    # 左卡片：GitHub
    add_rect(slide, Inches(0.5), res_y, res_w, res_h, COL_SLATE_800)
    add_rect(slide, Inches(0.5), res_y, Inches(0.08), res_h, COL_EMERALD_500)
    add_chip(slide, Inches(0.8), res_y + Inches(0.25), Inches(1.4), Inches(0.32),
             "GITHUB", fill_color=COL_EMERALD_500, font_size=10)
    add_text_box(slide, Inches(0.8), res_y + Inches(0.7), res_w - Inches(0.4), Inches(0.4),
                 "Auto-Empirical-Research-Skills", font_size=14, bold=True, color=COL_ZINC_100)
    add_text_box(slide, Inches(0.8), res_y + Inches(1.05), res_w - Inches(0.4), Inches(0.4),
                 "github.com/brycewang-stanford/Auto-Empirical-Research-Skills",
                 font_size=10, color=COL_INDIGO_500, font_name="Consolas")
    add_text_box(slide, Inches(0.8), res_y + Inches(1.35), res_w - Inches(0.4), Inches(0.4),
                 "Path: skills/69-Paper-WorkFlow · 配套 did_demo.ipynb 一键可跑",
                 font_size=10, color=COL_ZINC_400)

    # 右卡片：产品
    rx = Inches(6.75)
    add_rect(slide, rx, res_y, res_w, res_h, COL_SLATE_800)
    add_rect(slide, rx, res_y, Inches(0.08), res_h, COL_ROSE_500)
    add_chip(slide, rx + Inches(0.3), res_y + Inches(0.25), Inches(1.4), Inches(0.32),
             "PRODUCT", fill_color=COL_ROSE_500, font_size=10)
    add_text_box(slide, rx + Inches(0.3), res_y + Inches(0.7), res_w - Inches(0.4), Inches(0.4),
                 "CoPaper.AI · 论文协作平台",
                 font_size=14, bold=True, color=COL_ZINC_100)
    add_text_box(slide, rx + Inches(0.3), res_y + Inches(1.05), res_w - Inches(0.4), Inches(0.4),
                 "copaper.ai  ·  Stanford REAP 维护",
                 font_size=10, color=COL_ROSE_500, font_name="Consolas")
    add_text_box(slide, rx + Inches(0.3), res_y + Inches(1.35), res_w - Inches(0.4), Inches(0.4),
                 "面向实证研究的产学 AI 工具箱 · 与 Paper-WorkFlow 共同落地",
                 font_size=10, color=COL_ZINC_400)

    # 底部署名
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                 "Stanford REAP × CoPaper.AI · 实证研究 AI 助手 · MIT License",
                 font_size=10, color=COL_ZINC_500, align=PP_ALIGN.CENTER, font_name="Helvetica")


# ============================================================
# 主入口
# ============================================================
def main():
    print("🎨 开始生成 Paper-WorkFlow 介绍 PPT (5 页, 含开源项目矩阵)...")
    slide_1_open_source()
    print("  ✓ Slide 1: 开源项目矩阵 (Cover)")
    slide_2_overview()
    print("  ✓ Slide 2: 一图看懂")
    slide_3_case_intro()
    print("  ✓ Slide 3: 案例介绍 (CFPS 互联网)")
    slide_4_rigor()
    print("  ✓ Slide 4: 研究级标准")
    slide_5_call_to_action()
    print("  ✓ Slide 5: 触发语 + 资源")

    prs.save(OUTPUT)
    print(f"\n✅ PPT 已保存: {OUTPUT}")
    print(f"   总页数: {len(prs.slides)}")


if __name__ == "__main__":
    main()
