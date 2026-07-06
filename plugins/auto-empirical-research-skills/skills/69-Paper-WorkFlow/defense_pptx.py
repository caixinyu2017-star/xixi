#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
defense_pptx.py — 论文答辩 PPT 生成器（参数化版本）

与 build_pptx.py 的关系：
  - build_pptx.py：硬编码的 30 页 skill 演示版
  - defense_pptx.py（本文）：从工作区 + YAML/JSON 配置自动生成学位论文/期刊论文
    答辩 PPT，按 [`references/chinese-journals.md`](../references/chinese-journals.md) §6.5
    的标准结构。

用法：
    # 1) 论文答辩 PPT（默认结构）
    python3 defense_pptx.py --workspace <paper_workspace_dir> --type thesis

    # 2) 期刊汇报 PPT（顶刊汇报会）
    python3 defense_pptx.py --workspace <paper_workspace_dir> --type journal-talk

    # 3) 显式提供配置文件
    python3 defense_pptx.py --config <config.yaml> --output 答辩.pptx

依赖：python-pptx（pip install python-pptx）
"""

from __future__ import annotations

import argparse
import json
import os
import re
import sys
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Any

try:
    import yaml  # type: ignore
except ImportError:
    yaml = None  # YAML 可选；只支持 JSON

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:
    print("缺少 python-pptx，请先 pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ============================================================================
# 主题色（与 build_pptx.py 保持一致，跨 PPT 系列视觉统一）
# ============================================================================
NAVY    = RGBColor(0x1F, 0x3A, 0x5F)
NAVY_DK = RGBColor(0x14, 0x28, 0x44)
TEAL    = RGBColor(0x2E, 0x86, 0x86)
ORANGE  = RGBColor(0xE0, 0x7B, 0x39)
GOLD    = RGBColor(0xD9, 0xA5, 0x21)
INK     = RGBColor(0x22, 0x26, 0x33)
GREY    = RGBColor(0x5B, 0x63, 0x72)
LGREY   = RGBColor(0xEC, 0xEF, 0xF3)
CARD    = RGBColor(0xF5, 0xF7, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LINE    = RGBColor(0xD3, 0xDA, 0xE4)
GREEN   = RGBColor(0x27, 0xAE, 0x60)
RED     = RGBColor(0xC0, 0x39, 0x2B)

CJK_FONT = "Microsoft YaHei"
EN_FONT  = "Calibri"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


# ============================================================================
# 数据结构：答辩 PPT 配置
# ============================================================================
@dataclass
class SlideSpec:
    """单张幻灯片的规格。

    kind 决定该页的渲染模板（见 _render_* 系列函数）：
      - cover         ：封面
      - agenda        ：目录 / 路标
      - section_divider：章节分割页
      - content_text  ：纯文字要点页
      - content_bullets：项目符号要点页
      - finding       ：核心研究发现页（高亮设计）
      - table_view    ：结果表格页
      - figure_view   ：结果图页
      - thank_you     ：致谢 / 答辩结束
    """
    kind: str
    title: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    notes: str = ""          # 备注 / 演讲者注释（不进入 PPT 主体，写入 notes）
    accent: str = "navy"     # navy/teal/orange/gold/green/red
    footer: str = ""         # 论文 / 研究短名

    # 可选外部资源
    table_csv: str = ""      # 表格数据（CSV 路径）
    image_path: str = ""     # 图（PNG/PDF 路径）
    two_column: bool = False  # 内容页是否左右分栏
    left_bullets: list[str] = field(default_factory=list)
    right_bullets: list[str] = field(default_factory=list)


@dataclass
class DefenseConfig:
    """答辩 PPT 整体配置。"""
    title: str
    subtitle: str = ""
    presenter: str = ""
    advisor: str = ""
    degree: str = "硕士"     # 本科 / 硕士 / 博士
    department: str = ""
    institution: str = ""
    date: str = ""
    duration_min: int = 15   # 答辩时长（分钟）
    type: str = "thesis"     # thesis / journal-talk

    # 自动从 duration_min 推算每页时间（不需要再手动配）
    @property
    def total_slides(self) -> int:
        # thesis 模板：10-12 节 ≈ 18-28 张
        # journal-talk 模板：8-10 节 ≈ 15-22 张
        return {"thesis": 22, "journal-talk": 18}.get(self.type, 22)

    # 演讲者备注（每页）
    notes: dict[str, str] = field(default_factory=dict)

    # 核心内容
    background: list[str] = field(default_factory=list)     # 选题背景
    significance: list[str] = field(default_factory=list)  # 研究意义
    literature: list[str] = field(default_factory=list)    # 文献综述要点
    theory_hypotheses: list[str] = field(default_factory=list)  # 理论与假设
    data_description: str = ""                              # 数据描述（短文）
    identification: str = ""                                # 识别策略（短文）
    main_findings: list[dict] = field(default_factory=list)  # 主要发现，每项 dict 含 title/key/result
    robustness: list[str] = field(default_factory=list)    # 稳健性
    mechanism: list[str] = field(default_factory=list)      # 机制与异质性
    contributions: list[str] = field(default_factory=list) # 贡献
    limitations: list[str] = field(default_factory=list)    # 局限与展望
    qna_anticipation: list[str] = field(default_factory=list)  # 预演 Q&A

    # 自动元数据
    @classmethod
    def from_workspace(cls, workspace: Path) -> "DefenseConfig":
        """从 paper_workspace/ 目录读取 evidence_ledger / main_results.json / design_register 等生成 config。"""
        ws = Path(workspace)
        intake = _read_md(ws / "00_meta" / "intake.md") or ""
        design_reg = _read_md(ws / "03_analysis" / "design_register.md") or ""
        sample_audit = _read_md(ws / "02_data" / "sample_audit.md") or ""
        results = _read_json(ws / "03_analysis" / "results" / "main_results.json") or {}
        summary = _read_md(ws / "03_analysis" / "results" / "summary.md") or ""
        proposal = _read_md(ws / "01_proposal" / "proposal.md") or ""

        # 提取 metadata
        title = _first_match(proposal, r"#\s+(.+?)$") or _first_match(intake, r"论文题目[：:]\s*(.+?)$") or "论文答辩"
        presenter = _first_match(intake, r"(?:学生|作者|presenter)[：:]\s*(.+?)$") or ""
        advisor = _first_match(intake, r"(?:导师|advisor)[：:]\s*(.+?)$") or ""
        institution = _first_match(intake, r"(?:学校|institution)[：:]\s*(.+?)$") or ""
        department = _first_match(intake, r"(?:学院|系|department)[：:]\s*(.+?)$") or ""
        degree = _first_match(intake, r"(?:学位|degree)[：:]\s*(本科|硕士|博士)") or "硕士"
        date = _first_match(intake, r"(?:答辩日期|date)[：:]\s*(.+?)$") or ""

        # 提取关键信息（启发式）
        background = _extract_section(proposal, r"##\s*(?:一、)?\s*选题背景[与与]?意义?", max_items=4)
        significance = _extract_section(proposal, r"##\s*(?:二、)?\s*研究意义", max_items=3)
        literature = _extract_section(proposal, r"##\s*(?:三、)?\s*文献综述", max_items=6)
        theory = _extract_section(proposal, r"##\s*(?:四、)?\s*(?:理论分析|研究假设|理论框架)", max_items=5)

        # 数据描述
        data_desc = _first_match(sample_audit, r"##\s*数据(?:描述|来源)[\s\S]+?##") or ""

        # 识别策略
        ident = _extract_section(design_reg, r"##\s*(?:识别假设?|Identification)", max_items=4)
        ident_str = "；".join(ident) if ident else ""

        # 主要发现
        main_findings = _parse_findings(results)

        # 稳健性 / 机制 / 贡献（从 design_risk_ledger / evidence_ledger 提取）
        risk = _read_md(ws / "03_analysis" / "design_risk_ledger.md") or ""
        evidence = _read_md(ws / "00_meta" / "evidence_ledger.md") or ""
        robustness = _extract_section(risk, r"##\s*稳健性", max_items=5)
        mechanism = _extract_section(summary, r"##\s*(?:机制|异质性)", max_items=4)
        contributions = _extract_section(evidence, r"##\s*主要发现|##\s*核心贡献", max_items=4)

        return cls(
            title=title,
            presenter=presenter,
            advisor=advisor,
            degree=degree,
            department=department,
            institution=institution,
            date=date,
            background=background or ["（请在 01_proposal/proposal.md 补充选题背景）"],
            significance=significance or ["（请在 proposal.md 补充研究意义）"],
            literature=literature or ["（请在 proposal.md 补充文献综述）"],
            theory_hypotheses=theory or ["（请在 proposal.md 补充理论分析与研究假设）"],
            data_description=data_desc or "（请在 02_data/sample_audit.md 补充数据描述）",
            identification=ident_str or "（请在 03_analysis/design_register.md 补充识别策略）",
            main_findings=main_findings,
            robustness=robustness or ["（请在 03_analysis/robustness 补充）"],
            mechanism=mechanism or ["（请在 summary.md 补充机制分析）"],
            contributions=contributions or ["（请在 evidence_ledger.md 补充）"],
        )


# ============================================================================
# 辅助函数
# ============================================================================
def _read_md(path: Path) -> str | None:
    try:
        return Path(path).read_text(encoding="utf-8")
    except FileNotFoundError:
        return None


def _read_json(path: Path) -> dict | None:
    try:
        return json.loads(Path(path).read_text(encoding="utf-8"))
    except (FileNotFoundError, json.JSONDecodeError):
        return None


def _first_match(text: str | None, pattern: str) -> str | None:
    if not text:
        return None
    m = re.search(pattern, text, re.MULTILINE)
    return m.group(1).strip() if m else None


def _extract_section(text: str | None, header_pattern: str, max_items: int = 5) -> list[str]:
    """从 Markdown 中提取 ## 标题后的项目符号列表（粗略启发式）。"""
    if not text:
        return []
    m = re.search(header_pattern + r"([\s\S]*?)(?=\n##\s|\Z)", text)
    if not m:
        return []
    body = m.group(1)
    bullets = re.findall(r"^[\s]*[-*+]\s+(.+?)$", body, re.MULTILINE)
    if not bullets:
        # 没找到项目符号，提取段落
        paragraphs = [p.strip() for p in body.split("\n\n") if p.strip() and not p.strip().startswith("#")]
        bullets = paragraphs
    return [b.strip() for b in bullets[:max_items] if b.strip()]


def _parse_findings(results: dict | None) -> list[dict]:
    """从 main_results.json 提取主要发现（每个 estimate 取一个 finding）。"""
    if not results:
        return []
    findings = []
    estimates = results.get("estimates") or results.get("main_results") or []
    if not isinstance(estimates, list):
        return []
    for est in estimates[:6]:  # 最多 6 个发现
        coef = est.get("coefficient") or est.get("estimate") or est.get("coef")
        se = est.get("std_error") or est.get("se")
        var = est.get("variable") or est.get("term") or est.get("name") or "?"
        if coef is None:
            continue
        pval = est.get("p_value") or est.get("pvalue") or ""
        sig = ""
        if isinstance(pval, (int, float)):
            sig = "显著" if pval < 0.05 else "不显著"
        elif isinstance(pval, str) and pval.endswith("*"):
            sig = "显著"
        result = est.get("result") or est.get("interpretation") or ""
        findings.append({
            "title": est.get("title") or f"{var} 效应",
            "key": f"β = {coef:.3f}" + (f" (SE = {se:.3f})" if se is not None else ""),
            "result": result or sig,
        })
    return findings


# ============================================================================
# PPT 渲染原语（参考 build_pptx.py 的视觉风格，保持一致）
# ============================================================================
def _set_cjk(run, font_name=CJK_FONT):
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    # 设置东亚字体
    for tag in ("eastAsia", "cs", "ascii", "hAnsi"):
        existing = rPr.find(qn("a:{}".format(tag)))
        if existing is not None:
            rPr.remove(existing)
        new = rPr.makeelement(qn("a:{}".format(tag)), {"typeface": font_name})
        rPr.append(new)


def bg(slide, color):
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg_shape.line.fill.background()
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    return bg_shape


def rect(slide, l, t, w, h, color, line_color=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        if line_w is not None:
            s.line.width = line_w
    return s


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    return tb, tf


def para(tf, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, space_after=6, font=CJK_FONT):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text and len(tf.paragraphs) == 1 else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    _set_cjk(r, font)
    return p


def color_map(accent: str) -> RGBColor:
    return {
        "navy": NAVY, "teal": TEAL, "orange": ORANGE, "gold": GOLD,
        "green": GREEN, "red": RED, "ink": INK, "grey": GREY,
    }.get(accent, NAVY)


# ============================================================================
# 页面模板
# ============================================================================
def _add_page_chrome(slide, page_no, total, footer_text=""):
    """每页底部：左下脚注，右下页码。"""
    if footer_text:
        tb, tf = textbox(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.35))
        para(tf, footer_text, size=10, color=GREY, align=PP_ALIGN.LEFT, space_after=0)
    tb, tf = textbox(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.35))
    para(tf, f"{page_no} / {total}", size=10, color=GREY, align=PP_ALIGN.RIGHT, space_after=0)


def _slide_cover(prs, cfg: DefenseConfig, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)

    # 装饰条
    rect(slide, 0, Inches(2.6), EMU_W, Inches(0.08), ORANGE)
    rect(slide, 0, Inches(4.6), EMU_W, Inches(0.04), GOLD)

    # 主标题
    tb, tf = textbox(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.4))
    para(tf, cfg.title, size=36, color=WHITE, bold=True, align=PP_ALIGN.LEFT, space_after=10)
    if cfg.subtitle:
        para(tf, cfg.subtitle, size=18, color=LGREY, align=PP_ALIGN.LEFT, space_after=4)

    # 副信息：学位 / 答辩人 / 导师 / 单位 / 日期
    info_lines = []
    if cfg.degree:
        info_lines.append(f"学位类型：{cfg.degree}学位论文答辩")
    if cfg.presenter:
        info_lines.append(f"答辩人：{cfg.presenter}")
    if cfg.advisor:
        info_lines.append(f"导师：{cfg.advisor}")
    if cfg.department:
        info_lines.append(f"院系：{cfg.department}")
    if cfg.institution:
        info_lines.append(f"单位：{cfg.institution}")
    if cfg.date:
        info_lines.append(f"日期：{cfg.date}")

    tb, tf = textbox(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2))
    for line in info_lines:
        para(tf, line, size=14, color=WHITE, align=PP_ALIGN.LEFT, space_after=4)

    # 右上角品牌
    tb, tf = textbox(slide, Inches(9.5), Inches(0.4), Inches(3.5), Inches(0.4))
    para(tf, "Paper-WorkFlow 答辩 PPT 自动生成", size=10, color=LGREY, align=PP_ALIGN.RIGHT, space_after=0)
    _add_page_chrome(slide, 1, total, footer_text=cfg.title[:30])


def _slide_agenda(prs, cfg: DefenseConfig, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)

    # 标题
    tb, tf = textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(1))
    para(tf, "目录", size=28, color=NAVY, bold=True, align=PP_ALIGN.LEFT, space_after=4)
    para(tf, f"答辩时长 {cfg.duration_min} 分钟", size=12, color=GREY, align=PP_ALIGN.LEFT, space_after=0)

    # 目录条目
    items = [
        ("01", "选题背景与研究意义", "2 张"),
        ("02", "文献综述", "1 张"),
        ("03", "理论框架与研究假设", "2 张"),
        ("04", "研究设计：数据 + 识别策略", "2 张"),
        ("05", "实证结果：主要发现", f"{min(len(cfg.main_findings), 6) or 5} 张"),
        ("06", "稳健性检验", f"{min(len(cfg.robustness), 3) or 2} 张"),
        ("07", "机制与异质性", f"{min(len(cfg.mechanism), 2) or 1} 张"),
        ("08", "结论与贡献", "1 张"),
        ("09", "研究局限与展望", "1 张"),
        ("10", "答辩致谢", "1 张"),
    ]
    tb, tf = textbox(slide, Inches(0.8), Inches(1.8), Inches(12), Inches(5.5))
    for code, name, n in items:
        p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = f"{code}  "
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = ORANGE
        r1.font.name = EN_FONT
        r2 = p.add_run()
        r2.text = name
        r2.font.size = Pt(16)
        r2.font.color.rgb = INK
        r2.font.name = CJK_FONT
        _set_cjk(r2, CJK_FONT)
        r3 = p.add_run()
        r3.text = f"    ({n})"
        r3.font.size = Pt(11)
        r3.font.color.rgb = GREY
        r3.font.name = CJK_FONT
        _set_cjk(r3, CJK_FONT)
    _add_page_chrome(slide, page_no, total, footer_text=cfg.title[:30])


def _slide_section_divider(prs, code: str, title: str, en: str, page_no, total, cfg: DefenseConfig):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY_DK)

    # 大编号
    tb, tf = textbox(slide, Inches(0.8), Inches(2.0), Inches(4), Inches(3))
    para(tf, code, size=120, color=ORANGE, bold=True, align=PP_ALIGN.LEFT, space_after=0)

    # 标题
    tb, tf = textbox(slide, Inches(5), Inches(2.5), Inches(7.5), Inches(1.5))
    para(tf, title, size=32, color=WHITE, bold=True, align=PP_ALIGN.LEFT, space_after=4)
    para(tf, en, size=16, color=GOLD, align=PP_ALIGN.LEFT, space_after=0)

    # 装饰
    rect(slide, Inches(5), Inches(4.2), Inches(1.5), Inches(0.06), ORANGE)
    _add_page_chrome(slide, page_no, total, footer_text=cfg.title[:30])


def _slide_bullets(prs, title: str, bullets: list[str], page_no, total, cfg: DefenseConfig,
                   accent: str = "navy", two_column: bool = False,
                   left_bullets: list[str] = None, right_bullets: list[str] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)

    # 标题
    tb, tf = textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    para(tf, title, size=24, color=NAVY, bold=True, align=PP_ALIGN.LEFT, space_after=4)
    rect(slide, Inches(0.6), Inches(1.1), Inches(1.2), Inches(0.05), color_map(accent))

    if two_column and left_bullets and right_bullets:
        # 两栏布局
        tb, tf = textbox(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(5.5))
        for b in left_bullets:
            para(tf, f"• {b}", size=14, color=INK, space_after=8)
        tb, tf = textbox(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5.5))
        for b in right_bullets:
            para(tf, f"• {b}", size=14, color=INK, space_after=8)
    else:
        tb, tf = textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
        for b in bullets:
            p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
            p.space_after = Pt(10)
            r1 = p.add_run()
            r1.text = "• "
            r1.font.size = Pt(16)
            r1.font.color.rgb = color_map(accent)
            r1.font.bold = True
            r1.font.name = CJK_FONT
            _set_cjk(r1, CJK_FONT)
            r2 = p.add_run()
            r2.text = b
            r2.font.size = Pt(15)
            r2.font.color.rgb = INK
            r2.font.name = CJK_FONT
            _set_cjk(r2, CJK_FONT)
    _add_page_chrome(slide, page_no, total, footer_text=cfg.title[:30])


def _slide_finding(prs, title: str, finding: dict, page_no, total, cfg: DefenseConfig):
    """核心研究发现页：大字号 + 高亮。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)

    tb, tf = textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    para(tf, title, size=24, color=NAVY, bold=True, align=PP_ALIGN.LEFT, space_after=4)
    rect(slide, Inches(0.6), Inches(1.1), Inches(1.2), Inches(0.05), ORANGE)

    # 副标题
    if finding.get("title"):
        tb, tf = textbox(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.6))
        para(tf, finding["title"], size=18, color=INK, align=PP_ALIGN.LEFT, space_after=2)

    # 大字显示关键数字
    if finding.get("key"):
        tb, tf = textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(2))
        para(tf, finding["key"], size=44, color=ORANGE, bold=True, align=PP_ALIGN.LEFT, space_after=8)

    # 结果说明
    if finding.get("result"):
        tb, tf = textbox(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2))
        para(tf, finding["result"], size=15, color=INK, align=PP_ALIGN.LEFT, space_after=4)

    _add_page_chrome(slide, page_no, total, footer_text=cfg.title[:30])


def _slide_text(prs, title: str, body: str, page_no, total, cfg: DefenseConfig, accent: str = "navy"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, WHITE)
    tb, tf = textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    para(tf, title, size=24, color=NAVY, bold=True, align=PP_ALIGN.LEFT, space_after=4)
    rect(slide, Inches(0.6), Inches(1.1), Inches(1.2), Inches(0.05), color_map(accent))

    tb, tf = textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
    for para_text in body.split("\n\n"):
        para(tf, para_text.strip(), size=15, color=INK, space_after=10)
    _add_page_chrome(slide, page_no, total, footer_text=cfg.title[:30])


def _slide_thank_you(prs, cfg: DefenseConfig, page_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    rect(slide, 0, Inches(3.3), EMU_W, Inches(0.08), ORANGE)
    rect(slide, 0, Inches(3.9), EMU_W, Inches(0.04), GOLD)

    tb, tf = textbox(slide, Inches(0.6), Inches(2.5), Inches(12), Inches(1))
    para(tf, "敬请各位老师批评指正", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=8)
    para(tf, "Thank you for your questions and comments", size=18, color=LGREY, align=PP_ALIGN.CENTER, space_after=0)

    # 答辩人信息
    tb, tf = textbox(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(1.5))
    if cfg.presenter:
        para(tf, f"答辩人：{cfg.presenter}", size=18, color=WHITE, align=PP_ALIGN.CENTER, space_after=4)
    if cfg.advisor:
        para(tf, f"导师：{cfg.advisor}", size=14, color=LGREY, align=PP_ALIGN.CENTER, space_after=2)
    if cfg.institution:
        para(tf, cfg.institution, size=14, color=LGREY, align=PP_ALIGN.CENTER, space_after=0)
    if cfg.date:
        para(tf, cfg.date, size=14, color=LGREY, align=PP_ALIGN.CENTER, space_after=0)
    _add_page_chrome(slide, page_no, total, footer_text=cfg.title[:30])


# ============================================================================
# 模板组装（按 chinese-journals.md §6.5 标准结构）
# ============================================================================
def build_thesis_pptx(cfg: DefenseConfig, output_path: Path):
    """学位论文答辩 PPT 模板：22 张。"""
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    total = cfg.total_slides

    # 1. 封面
    _slide_cover(prs, cfg, total)

    # 2. 目录
    _slide_agenda(prs, cfg, 2, total)

    # 3. 章节分割：01 选题背景
    _slide_section_divider(prs, "01", "选题背景与研究意义", "BACKGROUND & SIGNIFICANCE", 3, total, cfg)

    # 4. 选题背景
    _slide_bullets(prs, "选题背景", cfg.background[:5], 4, total, cfg, accent="teal")

    # 5. 研究意义
    _slide_bullets(prs, "研究意义：理论 + 现实", cfg.significance[:4], 5, total, cfg, accent="orange")

    # 6. 章节分割：02 文献综述
    _slide_section_divider(prs, "02", "文献综述", "LITERATURE REVIEW", 6, total, cfg)

    # 7. 文献综述要点
    _slide_bullets(prs, "文献综述与本文定位", cfg.literature[:6], 7, total, cfg, accent="teal")

    # 8. 章节分割：03 理论框架
    _slide_section_divider(prs, "03", "理论框架与研究假设", "THEORY & HYPOTHESES", 8, total, cfg)

    # 9. 理论框架
    _slide_bullets(prs, "理论分析与研究假设", cfg.theory_hypotheses[:5], 9, total, cfg, accent="teal")

    # 10. 章节分割：04 研究设计
    _slide_section_divider(prs, "04", "研究设计", "RESEARCH DESIGN", 10, total, cfg)

    # 11. 数据描述
    _slide_text(prs, "数据与样本", cfg.data_description or "（请在 02_data/sample_audit.md 补充）", 11, total, cfg, accent="teal")

    # 12. 识别策略
    _slide_text(prs, "识别策略", cfg.identification or "（请在 03_analysis/design_register.md 补充）", 12, total, cfg, accent="orange")

    # 13. 章节分割：05 实证结果
    _slide_section_divider(prs, "05", "实证结果", "EMPIRICAL RESULTS", 13, total, cfg)

    # 14-19. 主要发现（每项 1 张）
    page = 14
    for f in (cfg.main_findings or [{"title": "（主发现未自动提取）", "key": "请编辑 03_analysis/results/main_results.json 或 evidence_ledger.md", "result": ""}])[:6]:
        _slide_finding(prs, "主要发现", f, page, total, cfg)
        page += 1

    # 20. 章节分割：06 稳健性
    if page < total - 1:
        _slide_section_divider(prs, "06", "稳健性检验", "ROBUSTNESS", page, total, cfg)
        page += 1
        _slide_bullets(prs, "稳健性检验", cfg.robustness[:4], page, total, cfg, accent="green")
        page += 1

    # 21. 机制与异质性（如果还有页）
    if page < total - 1:
        _slide_bullets(prs, "机制分析与异质性", cfg.mechanism[:4], page, total, cfg, accent="orange")
        page += 1

    # 22. 结论与贡献
    if page < total - 1:
        _slide_bullets(prs, "研究结论与主要贡献", cfg.contributions[:4], page, total, cfg, accent="navy")
        page += 1

    # 23. 局限与展望
    if page < total - 1:
        _slide_bullets(prs, "研究局限与未来展望", cfg.limitations[:3] or ["样本期可能限制外推；待更长时段验证", "机制识别可进一步引入实验设计", "政策评估可拓展至区域间比较"], page, total, cfg, accent="red")
        page += 1

    # 24. 答辩致谢
    _slide_thank_you(prs, cfg, page, total)

    # 实际生成页数
    actual_total = page
    prs.save(str(output_path))
    print(f"✓ 已生成 {output_path}（{actual_total} 页 / 模板估算 {total} 页）")
    return output_path


def build_journal_talk_pptx(cfg: DefenseConfig, output_path: Path):
    """期刊汇报 PPT 模板：18 张。"""
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    total = 18

    # 1. 封面
    _slide_cover(prs, cfg, total)

    # 2. 目录
    _slide_agenda(prs, cfg, 2, total)

    # 3. 选题背景
    _slide_section_divider(prs, "01", "Motivation", "选题背景与研究意义", 3, total, cfg)
    _slide_bullets(prs, "研究问题", cfg.background[:4], 4, total, cfg, accent="teal")

    # 4. 文献综述
    _slide_section_divider(prs, "02", "Related Literature", "文献综述", 5, total, cfg)
    _slide_bullets(prs, "文献定位", cfg.literature[:5], 6, total, cfg, accent="teal")

    # 5. 理论框架
    _slide_section_divider(prs, "03", "Framework", "理论框架", 7, total, cfg)
    _slide_text(prs, "理论框架", cfg.identification or "请补充", 8, total, cfg, accent="teal")

    # 6. 数据与识别
    _slide_section_divider(prs, "04", "Data & Identification", "数据与识别策略", 9, total, cfg)
    _slide_text(prs, "数据描述", cfg.data_description or "请补充", 10, total, cfg, accent="teal")

    # 7. 主要发现
    _slide_section_divider(prs, "05", "Main Results", "实证结果", 11, total, cfg)
    page = 12
    for f in (cfg.main_findings or [{"title": "（未自动提取）", "key": "请编辑 main_results.json", "result": ""}])[:5]:
        _slide_finding(prs, "Main Finding", f, page, total, cfg)
        page += 1

    # 8. 稳健性
    if page < total - 1:
        _slide_bullets(prs, "Robustness", cfg.robustness[:3], page, total, cfg, accent="green")
        page += 1

    # 9. 贡献
    if page < total - 1:
        _slide_bullets(prs, "Contribution", cfg.contributions[:4], page, total, cfg, accent="navy")
        page += 1

    # 致谢
    _slide_thank_you(prs, cfg, page, total)

    prs.save(str(output_path))
    print(f"✓ 已生成 {output_path}（{page} 页）")
    return output_path


# ============================================================================
# CLI
# ============================================================================
def main():
    p = argparse.ArgumentParser(description="答辩 PPT 自动生成（Paper-WorkFlow）")
    p.add_argument("--workspace", help="paper_workspace 目录路径")
    p.add_argument("--config", help="手动指定的 config.yaml 或 config.json 路径")
    p.add_argument("--output", default="答辩.pptx", help="输出 .pptx 路径")
    p.add_argument("--type", choices=["thesis", "journal-talk"], default="thesis")
    p.add_argument("--duration", type=int, default=15, help="答辩时长（分钟）")
    args = p.parse_args()

    # 加载 config
    if args.config:
        cfg_path = Path(args.config)
        raw = cfg_path.read_text(encoding="utf-8")
        if cfg_path.suffix in (".yaml", ".yml"):
            if yaml is None:
                print("需要 PyYAML 来解析 YAML", file=sys.stderr)
                sys.exit(1)
            data = yaml.safe_load(raw)
        else:
            data = json.loads(raw)
        cfg = DefenseConfig(**data)
    elif args.workspace:
        cfg = DefenseConfig.from_workspace(Path(args.workspace))
    else:
        print("必须提供 --workspace 或 --config", file=sys.stderr)
        sys.exit(1)

    cfg.type = args.type
    cfg.duration_min = args.duration
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if cfg.type == "thesis":
        build_thesis_pptx(cfg, output_path)
    else:
        build_journal_talk_pptx(cfg, output_path)


if __name__ == "__main__":
    main()
